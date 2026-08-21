from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import copy

G = RGBColor(27, 186, 141)
B = RGBColor(22, 71, 114)
W = RGBColor(255, 255, 255)
D = RGBColor(40, 40, 50)
GY = RGBColor(130, 130, 140)
LB = RGBColor(230, 242, 255)
LG = RGBColor(240, 248, 240)
R = RGBColor(220, 60, 60)
P = RGBColor(140, 60, 180)
A = RGBColor(230, 160, 30)
T = RGBColor(60, 130, 180)

MAIN = r"C:\Users\sumit\Downloads\Catchq\CatchQ_SeedFund_PitchDeck_final.pptx"
OUT = r"C:\Users\sumit\Downloads\Catchq\CatchQ_SeedFund_PitchDeck_v3.pptx"

prs = Presentation(MAIN)
SW = Inches(13.333)
SH = Inches(7.5)

def add_rect(slide, x, y, w, h, fc, lc=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fc
    if lc: s.line.color.rgb = lc; s.line.width = Pt(1)
    else: s.line.fill.background()
    return s

def add_rrect(slide, x, y, w, h, fc, lc=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fc
    if lc: s.line.color.rgb = lc; s.line.width = Pt(1)
    else: s.line.fill.background()
    return s

def add_circle(slide, x, y, sz, fc):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, sz, sz)
    s.fill.solid(); s.fill.fore_color.rgb = fc; s.line.fill.background()
    return s

def add_text(slide, x, y, w, h, text, sz=12, c=D, b=False, a=PP_ALIGN.LEFT, align=None):
    if align is not None: a = align
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz)
    p.font.color.rgb = c; p.font.bold = b; p.alignment = a; p.font.name = 'Calibri'
    return tb

def add_ml(slide, x, y, w, h, lines, sz=9, c=D, b=False, a=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(sz); p.font.color.rgb = c
        p.font.bold = b; p.alignment = a; p.font.name = 'Calibri'; p.space_after = Pt(3)
    return tb

def topbar(slide, title, subtitle):
    add_rect(slide, Inches(0), Inches(0), SW, Inches(1.2), B)
    add_rect(slide, Inches(0), Inches(1.2), SW, Inches(0.08), G)
    add_text(slide, Inches(0.5), Inches(0.15), Inches(10), Inches(0.7), title, 28, W, True)
    add_text(slide, Inches(0.5), Inches(0.75), Inches(10), Inches(0.4), subtitle, 13, RGBColor(170, 200, 230))

def botbar(slide, num, total=31):
    add_rect(slide, Inches(0), SH - Inches(0.35), SW, Inches(0.35), RGBColor(20, 55, 90))
    add_text(slide, Inches(0.5), SH - Inches(0.33), Inches(6), Inches(0.3),
        'CatchQ  |  IIT Kharagpur Platinum Jubilee Seed Fund Application', 8, RGBColor(140, 160, 180))
    add_text(slide, SW - Inches(1.5), SH - Inches(0.33), Inches(1.2), Inches(0.3),
        f'{num} / {total}', 8, RGBColor(140, 160, 180), align=PP_ALIGN.RIGHT)

def add_new_slide():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    for sh in list(sl.shapes):
        sp = sh._element; sp.getparent().remove(sp)
    return sl

# ================================================================
# SLIDE 1: COMPETITIVE LANDSCAPE
# ================================================================
sl = add_new_slide()
topbar(sl, 'Competitive Landscape', 'CatchQ vs alternatives - Only platform with queue + AI + bed management + WhatsApp')
botbar(sl, 1, 31)

# Comparison table
headers = ['Feature', 'CatchQ', 'Practo', 'ClinicEase', 'qTest', 'Paper/WhatsApp']
hcolors = [D, G, T, P, A, GY]
col_w = [Inches(2.2), Inches(2), Inches(2), Inches(2), Inches(2), Inches(2)]
x_start = Inches(0.6)

# Header row
x = x_start
for h, hc, w in zip(headers, hcolors, col_w):
    add_rect(sl, x, Inches(1.6), w, Inches(0.4), hc)
    add_text(sl, x, Inches(1.62), w, Inches(0.35), h, 10, W, True, PP_ALIGN.CENTER)
    x += w + Inches(0.05)

# Feature rows
features = [
    ('Real-time Queue', 'YES', 'NO', 'NO', 'NO', 'NO'),
    ('AI Scheduling', 'In Progress', 'NO', 'NO', 'NO', 'NO'),
    ('WhatsApp Integration', 'In Progress', 'Limited', 'NO', 'NO', 'YES'),
    ('Bed Management', 'YES', 'NO', 'NO', 'NO', 'NO'),
    ('Smart Billing', 'YES', 'YES', 'YES', 'NO', 'NO'),
    ('Digital Prescriptions', 'In Progress', 'YES', 'YES', 'YES', 'NO'),
    ('Patient App', 'YES', 'YES', 'NO', 'NO', 'NO'),
    ('Multi-Doctor Support', 'YES', 'YES', 'YES', 'YES', 'NO'),
    ('Analytics Dashboard', 'YES', 'YES', 'Basic', 'NO', 'NO'),
    ('Offline Support', 'In Progress', 'NO', 'NO', 'NO', 'YES'),
    ('Price/month', 'Rs.2,999', 'Rs.5,000+', 'Rs.4,000+', 'Rs.3,500', 'FREE'),
    ('Setup Time', '10 min', '2-3 days', '1 week', '1 week', '0'),
]

y = Inches(2.05)
for feat, cq, pr, ce, qt, pw in features:
    x = x_start
    vals = [feat, cq, pr, ce, qt, pw]
    for i, (val, w) in enumerate(zip(vals, col_w)):
        fc = LG if val == 'YES' else (RGBColor(255, 240, 240) if val == 'NO' else W)
        if i == 1 and val == 'YES': fc = RGBColor(230, 255, 230)
        add_rect(sl, x, y, w, Inches(0.35), fc)
        txt_c = G if (i == 1 and val == 'YES') else (R if val == 'NO' else D)
        add_text(sl, x, y + Inches(0.02), w, Inches(0.3), val, 8, txt_c, i == 1, PP_ALIGN.CENTER)
        x += w + Inches(0.05)
    y += Inches(0.36)

# Key differentiators
add_rrect(sl, Inches(0.4), Inches(6.5), Inches(12.5), Inches(0.6), LG, G)
add_text(sl, Inches(0.6), Inches(6.52), Inches(12), Inches(0.55),
    'CATCHQ ADVANTAGE:  Only all-in-one platform (Queue + Beds + Billing)  |  WhatsApp-native (In Progress)  |  AI Roadmap  |  Rs. 2,999/mo (60% cheaper)  |  10-min setup vs weeks',
    9, G, True, PP_ALIGN.LEFT)

# ================================================================
# SLIDE 2: USE OF FUNDS BREAKDOWN
# ================================================================
sl = add_new_slide()
topbar(sl, 'Use of Funds', 'Detailed allocation of Rs. 30-35 Lakhs seed fund - every rupee mapped to growth')
botbar(sl, 2, 31)

# Total fund
add_rrect(sl, Inches(5), Inches(1.5), Inches(3.3), Inches(0.9), B)
add_text(sl, Inches(5), Inches(1.55), Inches(3.3), Inches(0.4), 'TOTAL SEED FUND', 12, RGBColor(170, 200, 230), align=PP_ALIGN.CENTER)
add_text(sl, Inches(5), Inches(1.9), Inches(3.3), Inches(0.4), 'Rs. 30-35 Lakhs', 22, W, True, PP_ALIGN.CENTER)

# Allocation cards
allocations = [
    ('Product Development', '40%', 'Rs. 12-14L', G, [
        'Backend engineer (2): Rs. 8L',
        'Frontend/Android (1): Rs. 4L',
        'Cloud infra (AWS): Rs. 1.5L',
        'AI/ML development: Rs. 1L',
    ]),
    ('Sales & Marketing', '25%', 'Rs. 7.5-9L', T, [
        'WhatsApp marketing: Rs. 3L',
        'Sales team (2): Rs. 4L',
        'Content & branding: Rs. 1L',
        'Events & demos: Rs. 1L',
    ]),
    ('Operations', '15%', 'Rs. 4.5-5L', A, [
        'Office space: Rs. 2L',
        'Legal & compliance: Rs. 0.5L',
        'Accounting: Rs. 0.5L',
        'Misc ops: Rs. 1L',
    ]),
    ('Contingency', '10%', 'Rs. 3-3.5L', P, [
        'Emergency fund: Rs. 2L',
        'Opportunity fund: Rs. 1L',
        'Market pivot: Rs. 0.5L',
        '',
    ]),
    ('Working Capital', '10%', 'Rs. 3-3.5L', R, [
        'Cash buffer: Rs. 2L',
        'Vendor payments: Rs. 1L',
        'Insurance: Rs. 0.5L',
        '',
    ]),
]

x = Inches(0.4)
for name, pct, amt, color, items in allocations:
    add_rrect(sl, x, Inches(2.7), Inches(2.4), Inches(4), W, color)
    add_rect(sl, x, Inches(2.7), Inches(2.4), Inches(0.5), color)
    add_text(sl, x, Inches(2.72), Inches(2.4), Inches(0.22), name, 10, W, True, PP_ALIGN.CENTER)
    add_text(sl, x, Inches(2.95), Inches(2.4), Inches(0.2), pct, 14, W, True, PP_ALIGN.CENTER)
    add_text(sl, x, Inches(3.35), Inches(2.4), Inches(0.3), amt, 11, color, True, PP_ALIGN.CENTER)
    add_rect(sl, x + Inches(0.2), Inches(3.7), Inches(2), Inches(0.03), color)
    y = Inches(3.85)
    for item in items:
        if item:
            add_text(sl, x + Inches(0.2), y, Inches(2), Inches(0.25), f"  {item}", 8, D)
        y += Inches(0.28)
    x += Inches(2.55)

# ================================================================
# SLIDE 3: GO-TO-MARKET STRATEGY
# ================================================================
sl = add_new_slide()
topbar(sl, 'Go-to-Market Strategy', 'WhatsApp-first acquisition targeting 50 clinics in Year 1')
botbar(sl, 3, 31)

# Phase cards
phases = [
    ('PHASE 1: Month 1-3', 'Foundation', B, [
        'Launch with 5 beta clinics',
        'IIT KGP alumni network',
        'Free 30-day trial',
        'WhatsApp onboarding flow',
        'Collect NPS feedback',
        'Iterate on core features',
    ], 'TARGET: 10 Clinics'),
    ('PHASE 2: Month 4-6', 'Validation', G, [
        'Referral program (1 month free)',
        'IIT KGP incubator demo day',
        'Healthcare conferences',
        'Google Ads (clinic keywords)',
        'YouTube tutorials',
        'Case studies from beta',
    ], 'TARGET: 25 Clinics'),
    ('PHASE 3: Month 7-9', 'Scaling', T, [
        'Hire 2 sales executives',
        'Tier-2 city expansion',
        'WhatsApp broadcast campaigns',
        'Doctor association partnerships',
        'Insurance company tie-ups',
        'Government scheme integration',
    ], 'TARGET: 40 Clinics'),
    ('PHASE 4: Month 10-12', 'Growth', A, [
        'Channel partner program',
        'Clinic chain partnerships',
        'Media coverage (HealthTech)',
        'Series A preparation',
        '1000+ clinic waitlist',
        'Product market fit confirmed',
    ], 'TARGET: 50+ Clinics'),
]

x = Inches(0.3)
for title, sub, color, items, target in phases:
    add_rrect(sl, x, Inches(1.6), Inches(3.1), Inches(5.2), W, color)
    add_rect(sl, x, Inches(1.6), Inches(3.1), Inches(0.6), color)
    add_text(sl, x, Inches(1.62), Inches(3.1), Inches(0.3), title, 10, W, True, PP_ALIGN.CENTER)
    add_text(sl, x, Inches(1.9), Inches(3.1), Inches(0.25), sub, 11, W, align=PP_ALIGN.CENTER)
    y = Inches(2.4)
    for item in items:
        add_rrect(sl, x + Inches(0.15), y, Inches(2.8), Inches(0.38), LG)
        add_rect(sl, x + Inches(0.15), y, Inches(0.05), Inches(0.38), color)
        add_text(sl, x + Inches(0.3), y + Inches(0.05), Inches(2.5), Inches(0.28), item, 8, D)
        y += Inches(0.42)
    # Target
    add_rrect(sl, x + Inches(0.3), Inches(5.8), Inches(2.5), Inches(0.45), color)
    add_text(sl, x + Inches(0.3), Inches(5.82), Inches(2.5), Inches(0.4), target, 11, W, True, PP_ALIGN.CENTER)
    x += Inches(3.3)

# Key channels
add_rrect(sl, Inches(0.4), Inches(6.9), Inches(12.5), Inches(0.4), LG, G)
add_text(sl, Inches(0.6), Inches(6.92), Inches(12), Inches(0.35),
    'KEY CHANNELS:  WhatsApp Business API  |  Google Ads  |  IIT KGP Network  |  Doctor Referrals  |  Healthcare Events  |  YouTube',
    9, G, True, PP_ALIGN.LEFT)

# ================================================================
# SLIDE 4: UNIT ECONOMICS
# ================================================================
sl = add_new_slide()
topbar(sl, 'Unit Economics', 'Sustainable SaaS metrics with strong LTV/CAC ratio and healthy margins')
botbar(sl, 4, 31)

# Key metrics row
metrics = [
    ('Customer Acquisition\nCost (CAC)', 'Rs. 5,000', 'Blended', G),
    ('Lifetime Value\n(LTV)', 'Rs. 1,80,000', '36-month', B),
    ('LTV / CAC\nRatio', '36x', 'Target: >3x', A),
    ('Payback\nPeriod', '2 months', 'Target: <12', T),
    ('Monthly Churn\nRate', '3.5%', 'Target: <5%', P),
    ('Gross\nMargin', '82%', 'Target: >70%', R),
]

x = Inches(0.4)
for name, val, note, color in metrics:
    add_rrect(sl, x, Inches(1.6), Inches(2), Inches(1.6), W, color)
    add_rect(sl, x, Inches(1.6), Inches(2), Inches(0.04), color)
    add_text(sl, x, Inches(1.7), Inches(2), Inches(0.5), name, 9, GY, align=PP_ALIGN.CENTER)
    add_text(sl, x, Inches(2.2), Inches(2), Inches(0.5), val, 18, color, True, PP_ALIGN.CENTER)
    add_text(sl, x, Inches(2.7), Inches(2), Inches(0.3), note, 8, GY, align=PP_ALIGN.CENTER)
    x += Inches(2.1)

# Revenue breakdown
add_text(sl, Inches(0.5), Inches(3.5), Inches(4), Inches(0.4), 'REVENUE MODEL', 13, B, True)
add_rect(sl, Inches(0.5), Inches(3.85), Inches(1.5), Inches(0.04), G)

rev_items = [
    ('SaaS Subscription', '70%', 'Rs. 2,999-9,999/mo per clinic', G),
    ('Transaction Fees', '15%', '2% on billing through platform', T),
    ('Premium Features', '10%', 'AI analytics, advanced reports', A),
    ('API & Integrations', '5%', 'Third-party plugin marketplace', P),
]

y = Inches(4.0)
for name, pct, desc, color in rev_items:
    add_rrect(sl, Inches(0.5), y, Inches(5.5), Inches(0.55), W, RGBColor(230, 230, 235))
    add_rect(sl, Inches(0.5), y, Inches(0.06), Inches(0.55), color)
    # Percentage badge
    add_rrect(sl, Inches(0.7), y + Inches(0.08), Inches(0.8), Inches(0.35), color)
    add_text(sl, Inches(0.7), y + Inches(0.1), Inches(0.8), Inches(0.3), pct, 10, W, True, PP_ALIGN.CENTER)
    add_text(sl, Inches(1.7), y + Inches(0.03), Inches(2.5), Inches(0.25), name, 10, D, True)
    add_text(sl, Inches(1.7), y + Inches(0.28), Inches(4), Inches(0.2), desc, 8, GY)
    y += Inches(0.6)

# Right side: CAC breakdown
add_text(sl, Inches(7), Inches(3.5), Inches(4), Inches(0.4), 'CAC BREAKDOWN', 13, B, True)
add_rect(sl, Inches(7), Inches(3.85), Inches(1.5), Inches(0.04), G)

cac_items = [
    ('WhatsApp Ads', 'Rs. 1,500', '30%', T),
    ('Google Ads', 'Rs. 1,200', '24%', B),
    ('Sales Team', 'Rs. 1,000', '20%', G),
    ('Content Marketing', 'Rs. 800', '16%', A),
    ('Events/Demos', 'Rs. 500', '10%', P),
]

y = Inches(4.0)
for name, amt, pct, color in cac_items:
    add_rrect(sl, Inches(7), y, Inches(5.5), Inches(0.5), W, RGBColor(230, 230, 235))
    add_rect(sl, Inches(7), y, Inches(0.06), Inches(0.5), color)
    # Progress bar
    bar_w = float(pct.strip('%')) / 100 * 3
    add_rrect(sl, Inches(9.5), y + Inches(0.1), Inches(bar_w), Inches(0.3), color)
    add_text(sl, Inches(7.2), y + Inches(0.03), Inches(2.2), Inches(0.2), name, 9, D, True)
    add_text(sl, Inches(7.2), y + Inches(0.25), Inches(2), Inches(0.2), amt, 8, color, True)
    add_text(sl, Inches(9.5) + Inches(bar_w) + Inches(0.1), y + Inches(0.1), Inches(0.8), Inches(0.3), pct, 8, color, True)
    y += Inches(0.55)

# Bottom note
add_rrect(sl, Inches(0.4), Inches(6.7), Inches(12.5), Inches(0.5), RGBColor(255, 248, 230), A)
add_text(sl, Inches(0.6), Inches(6.72), Inches(12), Inches(0.45),
    'KEY INSIGHT: WhatsApp-first acquisition projected to cut CAC by 60%  |  Referral program target: 40% of new clinics  |  82% gross margin enables aggressive growth',
    9, A, True, PP_ALIGN.LEFT)

# ================================================================
# SLIDE 5: ROADMAP / MILESTONES
# ================================================================
sl = add_new_slide()
topbar(sl, 'Roadmap & Milestones', '18-month execution plan tied to seed fund deployment')
botbar(sl, 5, 31)

# Timeline
add_rect(sl, Inches(0.5), Inches(3.8), Inches(12.3), Inches(0.08), B)

# Milestone markers
milestones = [
    ('Month 1-3', 'MVP Launch', B, Inches(0.5), [
        'Core queue management live',
        '5 beta clinics onboarded',
        'WhatsApp integration V1',
        'Billing module complete',
        'Patient app launched',
    ]),
    ('Month 4-6', 'Product-Market Fit', G, Inches(3.3), [
        'AI scheduling V1',
        '100+ clinics in pipeline',
        'NPS score > 50',
        'Rs. 3L MRR',
        '2000+ patients/month',
    ]),
    ('Month 7-9', 'Growth Engine', T, Inches(6.1), [
        'Bed management launched',
        '25 active clinics',
        'Referral program live',
        'Rs. 8L MRR',
        '10,000+ patients/month',
    ]),
    ('Month 10-12', 'Scale', A, Inches(8.9), [
        '50 paying clinics',
        'AI diagnostics V1',
        'Rs. 15L MRR',
        'Series A readiness',
        '25,000+ patients/month',
    ]),
    ('Month 13-18', 'Expansion', P, Inches(11.7), [
        '100+ clinics',
        'Tier-2 city launch',
        'Rs. 30L MRR',
        'Insurance integrations',
        'Government partnerships',
    ]),
]

for title, sub, color, x, items in milestones:
    # Marker
    add_circle(sl, x + Inches(0.4), Inches(3.65), Inches(0.35), color)
    add_text(sl, x + Inches(0.4), Inches(3.7), Inches(0.35), Inches(0.25), '', 10, W, True, PP_ALIGN.CENTER)
    # Card
    add_rrect(sl, x, Inches(1.5), Inches(2.6), Inches(2), W, color)
    add_rect(sl, x, Inches(1.5), Inches(2.6), Inches(0.4), color)
    add_text(sl, x, Inches(1.52), Inches(2.6), Inches(0.2), title, 9, W, True, PP_ALIGN.CENTER)
    add_text(sl, x, Inches(1.72), Inches(2.6), Inches(0.18), sub, 10, W, align=PP_ALIGN.CENTER)
    y = Inches(2.0)
    for item in items:
        add_text(sl, x + Inches(0.1), y, Inches(2.4), Inches(0.18), f"  {item}", 7.5, D)
        y += Inches(0.2)
    # Bottom card
    add_rrect(sl, x, Inches(4.3), Inches(2.6), Inches(2.4), LG, color)
    add_text(sl, x, Inches(4.4), Inches(2.6), Inches(0.25), 'KEY METRICS', 8, color, True, PP_ALIGN.CENTER)
    y = Inches(4.7)
    # Metrics for this phase
    phase_metrics = {
        0: ['5 clinics live', 'Rs. 1.5L MRR', '80% activation'],
        1: ['15 clinics', 'Rs. 3L MRR', 'NPS > 50'],
        2: ['25 clinics', 'Rs. 8L MRR', '40% referrals'],
        3: ['50 clinics', 'Rs. 15L MRR', 'Series A ready'],
        4: ['100+ clinics', 'Rs. 30L MRR', 'Tier-2 launch'],
    }
    phase_idx = milestones.index((title, sub, color, x, items))
    for m in phase_metrics.get(phase_idx, []):
        add_rrect(sl, x + Inches(0.2), y, Inches(2.2), Inches(0.35), W)
        add_text(sl, x + Inches(0.3), y + Inches(0.05), Inches(2), Inches(0.25), m, 8, D, align=PP_ALIGN.CENTER)
        y += Inches(0.4)

# Bottom note
add_rrect(sl, Inches(0.4), Inches(6.9), Inches(12.5), Inches(0.4), LG, G)
add_text(sl, Inches(0.6), Inches(6.92), Inches(12), Inches(0.35),
    'MILESTONES TIED TO FUNDING:  Rs. 30-35L covers 18 months runway  |  Breakeven at Month 14  |  Series A at Month 12-15',
    9, G, True, PP_ALIGN.LEFT)

# ================================================================
# SLIDE 6: SOCIAL IMPACT
# ================================================================
sl = add_new_slide()
topbar(sl, 'Social Impact', 'Bridging the healthcare access gap for 500M+ Indians in tier-2/3 cities')
botbar(sl, 6, 31)

# Impact stats
stats = [
    ('500M+', 'Indians in tier-2/3\ncities lack quality\nhealthcare access', B),
    ('70%', 'Clinics still use\npaper registers\nand WhatsApp', G),
    ('3x', 'Average wait time\nin non-metro clinics\nvs metro hospitals', A),
    ('45 min', 'Average consultation\ndelay due to poor\nqueue management', R),
    ('10 Lakh+', 'Clinics in India\naddressable market', T),
]

x = Inches(0.4)
for val, desc, color in stats:
    add_rrect(sl, x, Inches(1.5), Inches(2.4), Inches(1.8), W, color)
    add_rect(sl, x, Inches(1.5), Inches(2.4), Inches(0.04), color)
    add_text(sl, x, Inches(1.6), Inches(2.4), Inches(0.6), val, 28, color, True, PP_ALIGN.CENTER)
    add_text(sl, x + Inches(0.1), Inches(2.3), Inches(2.2), Inches(0.9), desc, 9, D, align=PP_ALIGN.CENTER)
    x += Inches(2.55)

# Impact areas
add_text(sl, Inches(0.5), Inches(3.6), Inches(5), Inches(0.4), 'IMPACT AREAS', 13, B, True)
add_rect(sl, Inches(0.5), Inches(3.95), Inches(1.5), Inches(0.04), G)

impacts = [
    ('Healthcare Access', 'Makes quality clinic management affordable for small clinics in tier-2/3 cities', G),
    ('Digital India', 'Transforms paper-based clinics to digital-first, aligned with government vision', B),
    ('Patient Experience', 'Reduces wait times by 60%, eliminates paper prescriptions, enables digital payments', T),
    ('Doctor Efficiency', 'Automates queue, billing, scheduling - doctors see 30% more patients daily', A),
    ('Employment', 'Creates jobs: sales, support, development in tier-2 cities', P),
    ('Women Health', 'Special modules for maternal care, menstrual tracking, women-only clinics', R),
]

y = Inches(4.1)
for name, desc, color in impacts:
    add_rrect(sl, Inches(0.5), y, Inches(6), Inches(0.45), W, RGBColor(230, 230, 235))
    add_rect(sl, Inches(0.5), y, Inches(0.06), Inches(0.45), color)
    add_circle(sl, Inches(0.7), y + Inches(0.07), Inches(0.28), color)
    add_text(sl, Inches(0.7), y + Inches(0.09), Inches(0.28), Inches(0.2), name[0], 9, W, True, PP_ALIGN.CENTER)
    add_text(sl, Inches(1.1), y + Inches(0.05), Inches(2), Inches(0.2), name, 9, D, True)
    add_text(sl, Inches(1.1), y + Inches(0.23), Inches(5.2), Inches(0.2), desc, 7.5, GY)
    y += Inches(0.48)

# Right: UN SDGs alignment
add_text(sl, Inches(7.5), Inches(3.6), Inches(5), Inches(0.4), 'UN SDG ALIGNMENT', 13, B, True)
add_rect(sl, Inches(7.5), Inches(3.95), Inches(1.5), Inches(0.04), G)

sdgs = [
    ('SDG 3: Good Health', 'Ensure healthy lives and promote well-being for all at all ages', G),
    ('SDG 8: Decent Work', 'Promote sustained, inclusive economic growth and decent work', B),
    ('SDG 9: Innovation', 'Build resilient infrastructure, promote inclusive industrialization', T),
    ('SDG 10: Reduced Inequality', 'Reduce inequality within and among countries', A),
    ('SDG 11: Sustainable Cities', 'Make cities inclusive, safe, resilient and sustainable', P),
]

y = Inches(4.1)
for name, desc, color in sdgs:
    add_rrect(sl, Inches(7.5), y, Inches(5.5), Inches(0.55), W, RGBColor(230, 230, 235))
    add_rect(sl, Inches(7.5), y, Inches(0.06), Inches(0.55), color)
    add_text(sl, Inches(7.7), y + Inches(0.03), Inches(3), Inches(0.2), name, 9, color, True)
    add_text(sl, Inches(7.7), y + Inches(0.25), Inches(5), Inches(0.25), desc, 7.5, GY)
    y += Inches(0.6)

# ================================================================
# SLIDE 7: RISK ANALYSIS
# ================================================================
sl = add_new_slide()
topbar(sl, 'Risk Analysis', 'Key risks identified with mitigation strategies - showing maturity and preparedness')
botbar(sl, 7, 31)

risks = [
    ('Competition', 'HIGH', R, 'Practo, ClinicEase have funding & brand',
     'Differentiate on WhatsApp-first, price, and all-in-one. Win on tier-2/3 where they are weak.', 'Focus on underserved segments'),
    ('Customer Churn', 'MEDIUM', A, 'Clinics may not adopt digital tools',
     'WhatsApp onboarding (zero learning curve). Free trial. Dedicated support. Show ROI in 2 weeks.', 'Reduce churn to <3%/month'),
    ('Technical Scalability', 'MEDIUM', A, 'System must handle 1000+ concurrent clinics',
     'Cloud-native (AWS). Auto-scaling. Load testing. Modular architecture.', '99.9% uptime SLA'),
    ('Regulatory', 'LOW', G, 'Healthcare data privacy (ABDM compliance)',
     'Build ABDM-compliant from Day 1. Data encryption. Regular audits.', 'Full compliance in 6 months'),
    ('Team Risk', 'MEDIUM', A, 'Key person dependency in early stage',
     'Document everything. Cross-train team. Hire 2nd backend dev in Month 2.', 'No single point of failure'),
    ('Market Timing', 'LOW', G, 'Clinics may not be ready for digital',
      'WhatsApp is already used by 90% of clinics. We are digitizing existing behavior, not creating new.', 'Beta validated with pilot clinics'),
]

x = Inches(0.3)
y_start = Inches(1.5)
for i, (name, level, color, risk, mitigation, outcome) in enumerate(risks):
    col = i % 3
    row = i // 3
    cx = Inches(0.3) + col * Inches(4.3)
    cy = Inches(1.5) + row * Inches(2.7)

    add_rrect(sl, cx, cy, Inches(4.1), Inches(2.5), W, color)
    add_rect(sl, cx, cy, Inches(4.1), Inches(0.45), color)
    add_text(sl, cx, cy + Inches(0.02), Inches(2.5), Inches(0.2), name, 11, W, True)
    add_rrect(sl, cx + Inches(2.8), cy + Inches(0.07), Inches(1.1), Inches(0.3), W)
    add_text(sl, cx + Inches(2.8), cy + Inches(0.08), Inches(1.1), Inches(0.25), level, 8, color, True, PP_ALIGN.CENTER)

    add_text(sl, cx + Inches(0.1), cy + Inches(0.5), Inches(3.8), Inches(0.3), f'RISK: {risk}', 8, D)
    add_text(sl, cx + Inches(0.1), cy + Inches(0.9), Inches(3.8), Inches(0.3), f'MITIGATION: {mitigation}', 8, G)
    add_text(sl, cx + Inches(0.1), cy + Inches(1.5), Inches(3.8), Inches(0.3), f'OUTCOME: {outcome}', 8, T)

# Bottom note
add_rrect(sl, Inches(0.4), Inches(6.8), Inches(12.5), Inches(0.4), LG, G)
add_text(sl, Inches(0.6), Inches(6.82), Inches(12), Inches(0.35),
    'RISK MANAGEMENT:  Monthly review  |  Contingency fund (10%)  |  Agile pivots  |  Diverse revenue streams  |  Strong unit economics',
    9, G, True, PP_ALIGN.LEFT)

# ================================================================
# SLIDE 8: PARTNERSHIPS
# ================================================================
sl = add_new_slide()
topbar(sl, 'Partnerships & Integrations', 'Strategic partnerships to accelerate growth and credibility')
botbar(sl, 8, 31)

# Partnership categories
cats = [
    ('Government & Policy', B, [
        ('Ayushman Bharat (PMJAY)', 'Integrate with 50Cr+ beneficiary database', 'In Discussion'),
        ('ABDM (Ayushman Bharat Digital)', 'Compliant with national health stack', 'Building'),
        ('State Health Departments', 'Pilot in West Bengal, Odisha, Jharkhand', 'Planned'),
        ('IIT KGP Incubator', 'Mentorship, network, credibility', 'Active'),
    ]),
    ('Technology Partners', G, [
        ('WhatsApp Business API', 'Official tech partner for healthcare', 'Applied'),
        ('AWS Activate', '$10K credits for startups', 'Approved'),
        ('Razorpay', 'Payment gateway integration', 'Live'),
        ('Google Maps API', 'Clinic location services', 'Live'),
    ]),
    ('Healthcare Partners', T, [
        ('Practo (potential)', 'API integration for referrals', 'Planned'),
        ('PharmEasy', 'Medicine delivery integration', 'In Discussion'),
        ('Thyrocare', 'Lab test booking integration', 'In Discussion'),
        ('1mg', 'Pharmacy integration', 'Planned'),
    ]),
    ('Channel Partners', A, [
        ('Doctor Associations', 'IMA, IAP state chapters', 'In Discussion'),
        ('Medical Equip Vendors', 'Bundle deals for clinics', 'Planned'),
        ('Insurance Companies', 'Cashless claim processing', 'Planned'),
        ('Chartered Accountants', 'Referral network for billing', 'Planned'),
    ]),
]

x = Inches(0.3)
for cat_name, color, partners in cats:
    add_rrect(sl, x, Inches(1.5), Inches(3.1), Inches(5.2), W, color)
    add_rect(sl, x, Inches(1.5), Inches(3.1), Inches(0.45), color)
    add_text(sl, x, Inches(1.52), Inches(3.1), Inches(0.4), cat_name, 11, W, True, PP_ALIGN.CENTER)

    y = Inches(2.1)
    for name, desc, status in partners:
        add_rrect(sl, x + Inches(0.1), y, Inches(2.9), Inches(1.0), LG)
        add_text(sl, x + Inches(0.2), y + Inches(0.05), Inches(2.7), Inches(0.2), name, 8, D, True)
        add_text(sl, x + Inches(0.2), y + Inches(0.28), Inches(2.7), Inches(0.35), desc, 7, GY)
        # Status badge
        sc = G if status == 'Live' else (T if status == 'Active' else (A if 'Discuss' in status or 'Applied' in status or 'Approved' in status else GY))
        add_rrect(sl, x + Inches(0.2), y + Inches(0.7), Inches(1.2), Inches(0.22), sc)
        add_text(sl, x + Inches(0.2), y + Inches(0.71), Inches(1.2), Inches(0.2), status, 6.5, W, True, PP_ALIGN.CENTER)
        y += Inches(1.05)
    x += Inches(3.3)

# ================================================================
# SLIDE 9: SCALABILITY PLAN
# ================================================================
sl = add_new_slide()
topbar(sl, 'Scalability Plan', 'From 50 to 1,000+ clinics - proven playbook with clear metrics')
botbar(sl, 9, 31)

# Growth stages
stages = [
    ('STAGE 1\nSeed', 'Month 1-6', '50 Clinics', 'Rs. 36L ARR', B, [
        'IIT KGP network',
        'WhatsApp direct outreach',
        'Free trial conversion',
        'Manual onboarding',
        'Founder-led sales',
    ]),
    ('STAGE 2\nSeries A', 'Month 7-18', '200 Clinics', 'Rs. 1.8 Cr ARR', G, [
        'Sales team (5 people)',
        'Channel partners',
        'Automated onboarding',
        'Referral program',
        'Content marketing',
    ]),
    ('STAGE 3\nSeries B', 'Month 19-36', '1,000 Clinics', 'Rs. 9.6 Cr ARR', T, [
        'National expansion',
        'Enterprise chains',
        'API marketplace',
        'International (SE Asia)',
        'AI diagnostics',
    ]),
    ('STAGE 4\nIPO', 'Year 3-5', '5,000+ Clinics', 'Rs. 50 Cr ARR', A, [
        'Platform ecosystem',
        'Insurance integration',
        'Health data marketplace',
        'Govt partnerships',
        'Global expansion',
    ]),
]

x = Inches(0.3)
for title, timeline, clinics, arr, color, items in stages:
    add_rrect(sl, x, Inches(1.5), Inches(3.1), Inches(5.2), W, color)
    add_rect(sl, x, Inches(1.5), Inches(3.1), Inches(0.7), color)
    add_text(sl, x, Inches(1.52), Inches(3.1), Inches(0.35), title, 11, W, True, PP_ALIGN.CENTER)
    add_text(sl, x, Inches(1.9), Inches(3.1), Inches(0.25), timeline, 9, W, align=PP_ALIGN.CENTER)

    # Metrics
    add_rrect(sl, x + Inches(0.2), Inches(2.4), Inches(2.7), Inches(0.8), LG)
    add_text(sl, x + Inches(0.3), Inches(2.45), Inches(2.5), Inches(0.3), clinics, 18, color, True, PP_ALIGN.CENTER)
    add_text(sl, x + Inches(0.3), Inches(2.8), Inches(2.5), Inches(0.25), arr, 11, D, True, PP_ALIGN.CENTER)
    add_text(sl, x + Inches(0.3), Inches(3.0), Inches(2.5), Inches(0.2), 'Annual Recurring Revenue', 7, GY, align=PP_ALIGN.CENTER)

    # Items
    y = Inches(3.4)
    for item in items:
        add_rrect(sl, x + Inches(0.15), y, Inches(2.8), Inches(0.35), LG)
        add_rect(sl, x + Inches(0.15), y, Inches(0.05), Inches(0.35), color)
        add_text(sl, x + Inches(0.3), y + Inches(0.05), Inches(2.5), Inches(0.25), item, 8, D)
        y += Inches(0.4)

    # Key metric
    add_rrect(sl, x + Inches(0.3), Inches(5.6), Inches(2.5), Inches(0.45), color)
    key_metrics = ['50 clinics', '200 clinics', '1,000 clinics', '5,000+ clinics']
    km = key_metrics[stages.index((title, timeline, clinics, arr, color, items))]
    add_text(sl, x + Inches(0.3), Inches(5.62), Inches(2.5), Inches(0.4), km, 13, W, True, PP_ALIGN.CENTER)
    x += Inches(3.3)

# Bottom
add_rrect(sl, Inches(0.4), Inches(6.8), Inches(12.5), Inches(0.4), LG, G)
add_text(sl, Inches(0.6), Inches(6.82), Inches(12), Inches(0.35),
    'SCALABILITY:  WhatsApp-first = zero marginal cost  |  Modular SaaS = easy expansion  |  Data network effects  |  Platform ecosystem at scale',
    9, G, True, PP_ALIGN.LEFT)

# ================================================================
# SAVE
# ================================================================
prs.save(OUT)
print(f"Final deck: {len(prs.slides)} slides")
print(f"Saved: {OUT}")
