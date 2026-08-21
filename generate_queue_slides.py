from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Brand colors
G = RGBColor(27, 186, 141)
B = RGBColor(22, 71, 114)
W = RGBColor(255, 255, 255)
D = RGBColor(40, 40, 50)
GY = RGBColor(130, 130, 140)
LB = RGBColor(230, 242, 255)
LG = RGBColor(240, 248, 240)
R = RGBColor(220, 60, 60)
P = RGBColor(140, 60, 180)
A = RGBColor(230, 160, 30)
T = RGBColor(60, 130, 180)

IMG_DIR = r"C:\Users\sumit\Downloads\Catchq\images"
OUT = r"C:\Users\sumit\Downloads\Catchq\CatchQ_Queue_TAM_Slides.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = Inches(13.333)
SH = Inches(7.5)

def add_rect(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_circle(slide, x, y, size, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text(slide, x, y, w, h, text, size=12, color=D, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    p.font.name = 'Calibri'
    return txBox

def add_multiline(slide, x, y, w, h, lines, size=10, color=D, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = align
        p.font.name = 'Calibri'
        p.space_after = Pt(4)
    return txBox

def add_image_safe(slide, img_name, x, y, w, h):
    path = os.path.join(IMG_DIR, img_name)
    if os.path.exists(path):
        try:
            return slide.shapes.add_picture(path, x, y, w, h)
        except:
            pass
    return None

def topbar(slide, title, subtitle):
    add_rect(slide, Inches(0), Inches(0), SW, Inches(1.2), B)
    add_rect(slide, Inches(0), Inches(1.2), SW, Inches(0.08), G)
    add_text(slide, Inches(0.5), Inches(0.15), Inches(10), Inches(0.7), title, 28, W, True)
    add_text(slide, Inches(0.5), Inches(0.75), Inches(10), Inches(0.4), subtitle, 13, RGBColor(170, 200, 230))

def botbar(slide, num, total=20):
    add_rect(slide, Inches(0), SH - Inches(0.35), SW, Inches(0.35), RGBColor(20, 55, 90))
    add_text(slide, Inches(0.5), SH - Inches(0.33), Inches(6), Inches(0.3),
        'CatchQ  |  IIT Kharagpur Platinum Jubilee Seed Fund Application', 8, RGBColor(140, 160, 180))
    add_text(slide, SW - Inches(1.5), SH - Inches(0.33), Inches(1.2), Inches(0.3),
        f'{num} / {total}', 8, RGBColor(140, 160, 180), align=PP_ALIGN.RIGHT)

# ================================================================
# SLIDE 1: QUEUE MANAGEMENT - OVERVIEW
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, Inches(0), Inches(0), SW, Inches(0), W)
topbar(slide, 'Queue Management System', 'Real-time queue with live sync between Web App and Android App')
botbar(slide, 1, 20)

# Left: Receptionist Web App
add_rounded_rect(slide, Inches(0.4), Inches(1.6), Inches(6.2), Inches(5.2), LB)
add_rect(slide, Inches(0.4), Inches(1.6), Inches(6.2), Inches(0.5), B)
add_text(slide, Inches(0.6), Inches(1.65), Inches(5), Inches(0.4), 'RECEPTIONIST WEB APP (CatchQ Dashboard)', 12, W, True)

# Web app mockup
add_rounded_rect(slide, Inches(0.6), Inches(2.3), Inches(5.8), Inches(4.2), W, RGBColor(220, 220, 225))

# Status bar
add_rect(slide, Inches(0.6), Inches(2.3), Inches(5.8), Inches(0.3), RGBColor(245, 245, 248))
add_text(slide, Inches(0.8), Inches(2.32), Inches(3), Inches(0.25), 'CatchQ Dashboard - Queue Management', 7, GY)

# Queue table header
add_rect(slide, Inches(0.8), Inches(2.8), Inches(5.4), Inches(0.35), B)
headers = ['#', 'Patient', 'Doctor', 'Time', 'Status']
widths = [0.4, 1.4, 1.4, 1.0, 1.2]
x = Inches(0.8)
for h, w in zip(headers, widths):
    add_text(slide, x, Inches(2.82), Inches(w), Inches(0.3), h, 8, W, True)
    x += Inches(w + 0.1)

# Queue rows with different statuses
patients = [
    ('1', 'Rahul Sharma', 'Dr. Patel (Cardio)', '10:00 AM', 'completed', G),
    ('2', 'Priya Singh', 'Dr. Patel (Cardio)', '10:15 AM', 'completed', G),
    ('3', 'Amit Kumar', 'Dr. Patel (Cardio)', '10:30 AM', 'serving', RGBColor(230, 160, 30)),
    ('4', 'Sneha Reddy', 'Dr. Patel (Cardio)', '10:45 AM', 'waiting', T),
    ('5', 'Vikram Joshi', 'Dr. Patel (Cardio)', '11:00 AM', 'waiting', T),
    ('6', 'Neha Gupta', 'Dr. Patel (Cardio)', '11:15 AM', 'waiting', T),
    ('7', 'Rajesh Verma', 'Dr. Patel (Cardio)', '11:30 AM', 'skipped', R),
    ('8', 'Kavita Nair', 'Dr. Patel (Cardio)', '11:45 AM', 'waiting', T),
    ('9', 'Sanjay Mehta', 'Dr. Patel (Cardio)', '12:00 PM', 'waiting', T),
    ('10', 'Pooja Desai', 'Dr. Patel (Cardio)', '12:15 PM', 'waiting', T),
]

y = Inches(3.2)
for num, patient, doctor, time, status, color in patients:
    if num in ['7']:
        # Highlighted skipped row
        add_rect(slide, Inches(0.8), y, Inches(5.4), Inches(0.32), RGBColor(255, 240, 240))
        add_rect(slide, Inches(0.8), y, Inches(5.4), Inches(0.02), R)
    x = Inches(0.8)
    vals = [num, patient, doctor, time, status.upper()]
    for val, w in zip(vals, widths):
        txt_color = color if val == status.upper() else D
        add_text(slide, x, y + Inches(0.02), Inches(w), Inches(0.28), val, 7, txt_color, val == status.upper())
        x += Inches(w + 0.1)
    y += Inches(0.33)

# Action buttons
add_rounded_rect(slide, Inches(0.8), Inches(6.2), Inches(1.2), Inches(0.3), G)
add_text(slide, Inches(0.8), Inches(6.22), Inches(1.2), Inches(0.26), 'NEXT PATIENT', 7, W, True, PP_ALIGN.CENTER)
add_rounded_rect(slide, Inches(2.1), Inches(6.2), Inches(1), Inches(0.3), R)
add_text(slide, Inches(2.1), Inches(6.22), Inches(1), Inches(0.26), 'SKIP', 7, W, True, PP_ALIGN.CENTER)
add_rounded_rect(slide, Inches(3.2), Inches(6.2), Inches(1), Inches(0.3), A)
add_text(slide, Inches(3.2), Inches(6.22), Inches(1), Inches(0.26), 'PAUSE', 7, W, True, PP_ALIGN.CENTER)

# Right: Android App
add_rounded_rect(slide, Inches(6.9), Inches(1.6), Inches(6.1), Inches(5.2), LG)
add_rect(slide, Inches(6.9), Inches(1.6), Inches(6.1), Inches(0.5), G)
add_text(slide, Inches(7.1), Inches(1.65), Inches(5), Inches(0.4), 'PATIENT ANDROID APP (CatchQ)', 12, W, True)

# Phone mockup
add_rounded_rect(slide, Inches(8.5), Inches(2.3), Inches(3), Inches(4.3), RGBColor(30, 30, 35))
add_rect(slide, Inches(8.6), Inches(2.5), Inches(2.8), Inches(3.8), W)

# Status bar
add_rect(slide, Inches(8.6), Inches(2.5), Inches(2.8), Inches(0.2), G)
add_text(slide, Inches(8.7), Inches(2.52), Inches(1), Inches(0.18), '9:41', 6, W, True)

# Queue info
add_text(slide, Inches(8.7), Inches(2.85), Inches(2.6), Inches(0.3), 'YOUR QUEUE STATUS', 8, B, True, PP_ALIGN.CENTER)
add_rect(slide, Inches(8.8), Inches(3.2), Inches(2.4), Inches(0.03), G)

# Big number
add_circle(slide, Inches(9.3), Inches(3.4), Inches(1.2), RGBColor(240, 248, 240))
add_text(slide, Inches(9.3), Inches(3.5), Inches(1.2), Inches(0.8), '#4', 28, G, True, PP_ALIGN.CENTER)
add_text(slide, Inches(8.7), Inches(4.7), Inches(2.6), Inches(0.25), 'Serving Now: #3', 9, D, True, PP_ALIGN.CENTER)
add_text(slide, Inches(8.7), Inches(5.0), Inches(2.6), Inches(0.25), 'Total in Queue: 8', 9, GY, align=PP_ALIGN.CENTER)
add_text(slide, Inches(8.7), Inches(5.3), Inches(2.6), Inches(0.25), 'Est. Wait: 15 min', 9, A, True, PP_ALIGN.CENTER)

# Appointment details
add_rounded_rect(slide, Inches(7.1), Inches(2.8), Inches(1.3), Inches(3.5), W, RGBColor(220, 220, 225))
add_text(slide, Inches(7.2), Inches(2.9), Inches(1.1), Inches(0.25), 'APPOINTMENT', 6, GY, True, PP_ALIGN.CENTER)
add_text(slide, Inches(7.2), Inches(3.2), Inches(1.1), Inches(0.2), 'Dr. Patel', 7, D, True, PP_ALIGN.CENTER)
add_text(slide, Inches(7.2), Inches(3.4), Inches(1.1), Inches(0.2), 'Cardiology', 6, GY, align=PP_ALIGN.CENTER)
add_text(slide, Inches(7.2), Inches(3.7), Inches(1.1), Inches(0.2), '10:45 AM', 7, B, True, PP_ALIGN.CENTER)
add_text(slide, Inches(7.2), Inches(4.0), Inches(1.1), Inches(0.2), 'Today', 6, G, align=PP_ALIGN.CENTER)

# Connection arrows
add_text(slide, Inches(6.3), Inches(3.5), Inches(0.6), Inches(0.4), '>>', 18, G, True, PP_ALIGN.CENTER)
add_text(slide, Inches(6.3), Inches(4.2), Inches(0.6), Inches(0.4), '<<', 18, G, True, PP_ALIGN.CENTER)

# Live sync label
add_rounded_rect(slide, Inches(7.1), Inches(5.8), Inches(5.9), Inches(0.5), RGBColor(255, 248, 230), A)
add_text(slide, Inches(7.3), Inches(5.82), Inches(5.5), Inches(0.45),
    'REAL-TIME SYNC: Status changes on web app instantly reflect on patient app via Socket.io WebSocket', 8, A, True, PP_ALIGN.LEFT)

# Bottom flow
add_rounded_rect(slide, Inches(0.4), Inches(6.9), Inches(12.5), Inches(0.4), LG, G)
add_text(slide, Inches(0.6), Inches(6.92), Inches(12), Inches(0.35),
    'Flow: Patient books via App/Web -> Receptionist sees in queue -> Status updates in real-time -> Patient sees live position on phone',
    9, G, True, PP_ALIGN.LEFT)

# ================================================================
# SLIDE 2: QUEUE STATUS FLOW
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, Inches(0), Inches(0), SW, Inches(0), W)
topbar(slide, 'Queue Status Flow', 'Patient journey through the queue system with real-time status updates')
botbar(slide, 2, 20)

# Flow diagram
statuses = [
    ('PENDING', 'Patient arrives\nor books online', RGBColor(200, 200, 210), 'Appointment created.\nAdded to queue.\nWaiting for check-in.', '12'),
    ('WAITING', 'In queue.\nPosition visible.', T, 'Queue position #N.\nEstimated wait time.\nReal-time updates.', '8'),
    ('SERVING', 'Doctor is\nseeing patient.', A, 'Currently with doctor.\nQueue advances.\nSocket.io live.', '1'),
    ('COMPLETED', 'Visit done.\nPatient leaves.', G, 'Receipt generated.\nFollow-up scheduled.\nQueue slot freed.', '2'),
    ('SKIPPED', 'Patient did\nnot arrive.', R, 'Auto-detected after\n15 min. Queue skips.\n+10 patients added.', '1'),
]

x = Inches(0.3)
for i, (status, desc, color, details, count) in enumerate(statuses):
    w = Inches(2.4)
    # Main card
    add_rounded_rect(slide, x, Inches(1.6), w, Inches(3.5), W, color)
    add_rect(slide, x, Inches(1.6), w, Inches(0.6), color)
    add_text(slide, x, Inches(1.65), w, Inches(0.5), status, 14, W, True, PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.1), Inches(2.3), w - Inches(0.2), Inches(0.6), desc, 10, D, align=PP_ALIGN.CENTER)
    # Count
    add_circle(slide, x + Inches(0.85), Inches(3.0), Inches(0.7), color)
    add_text(slide, x + Inches(0.85), Inches(3.1), Inches(0.7), Inches(0.5), count, 16, W, True, PP_ALIGN.CENTER)
    add_text(slide, x, Inches(3.7), w, Inches(0.25), 'patients', 8, GY, align=PP_ALIGN.CENTER)
    # Details
    add_multiline(slide, x + Inches(0.1), Inches(4.0), w - Inches(0.2), Inches(1), details.split('\n'), 8, GY, align=PP_ALIGN.CENTER)
    # Arrow
    if i < 4:
        add_text(slide, x + w - Inches(0.1), Inches(3.0), Inches(0.5), Inches(0.4), '>', 20, GY, True, PP_ALIGN.CENTER)
    x += w + Inches(0.2)

# Skip logic box
add_rounded_rect(slide, Inches(0.4), Inches(5.4), Inches(12.5), Inches(1.2), RGBColor(255, 245, 245), R)
add_rect(slide, Inches(0.4), Inches(5.4), Inches(0.1), Inches(1.2), R)
add_text(slide, Inches(0.7), Inches(5.5), Inches(3), Inches(0.4), 'SKIP LOGIC (Smart Queue)', 12, R, True)
add_multiline(slide, Inches(0.7), Inches(5.9), Inches(5.5), Inches(0.6), [
    'When a patient doesn\'t arrive within 15 minutes of their slot:',
    '1. System auto-marks as SKIPPED  2. Queue position shifts up',
    '3. Next +10 patients get updated positions  4. Push notification sent to affected patients',
], 9, D)

add_text(slide, Inches(7), Inches(5.5), Inches(3), Inches(0.4), 'WHY +10 PATIENTS?', 12, R, True)
add_multiline(slide, Inches(7), Inches(5.9), Inches(5.5), Inches(0.6), [
    'Adding buffer of 10 patients below current position ensures:',
    'Queue never appears empty  |  Reduces perceived wait time',
    'Doctor always has patients ready  |  No idle gaps between consultations',
], 9, D)

# ================================================================
# SLIDE 3: WEB vs APP COMPARISON
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, Inches(0), Inches(0), SW, Inches(0), W)
topbar(slide, 'Web App vs Android App', 'Dual interface - Receptionist controls queue, Patient views live status')
botbar(slide, 3, 20)

# Web App column
add_text(slide, Inches(0.5), Inches(1.6), Inches(6), Inches(0.4), 'RECEPTIONIST WEB APP', 14, B, True)
add_rect(slide, Inches(0.5), Inches(2.0), Inches(2), Inches(0.04), B)

web_features = [
    ('Queue Board', 'View all patients for all doctors in one dashboard', B),
    ('Status Control', 'Change status: Pending -> Waiting -> Serving -> Completed/Skipped', G),
    ('Skip Patient', 'One-click skip with auto-queue adjustment (+10 patients)', R),
    ('Doctor Filter', 'Filter queue by doctor, time slot, or date', T),
    ('Bulk Actions', 'Mark multiple patients as completed or skipped', A),
    ('Real-time Updates', 'Socket.io pushes changes instantly to patient apps', G),
    ('Analytics', 'Daily queue stats, avg wait time, no-show rate', P),
    ('Patient Search', 'Find patient by name, phone, or appointment ID', GY),
]

y = Inches(2.2)
for name, desc, color in web_features:
    add_rounded_rect(slide, Inches(0.5), y, Inches(6), Inches(0.5), W, RGBColor(230, 230, 235))
    add_rect(slide, Inches(0.5), y, Inches(0.06), Inches(0.5), color)
    add_circle(slide, Inches(0.7), y + Inches(0.08), Inches(0.3), color)
    add_text(slide, Inches(0.7), y + Inches(0.1), Inches(0.3), Inches(0.25), name[0], 10, W, True, PP_ALIGN.CENTER)
    add_text(slide, Inches(1.1), y + Inches(0.05), Inches(2), Inches(0.2), name, 9, D, True)
    add_text(slide, Inches(1.1), y + Inches(0.25), Inches(5), Inches(0.2), desc, 8, GY)
    y += Inches(0.55)

# Android App column
add_text(slide, Inches(6.8), Inches(1.6), Inches(6), Inches(0.4), 'PATIENT ANDROID APP', 14, G, True)
add_rect(slide, Inches(6.8), Inches(2.0), Inches(2), Inches(0.04), G)

app_features = [
    ('Queue Position', 'See your exact position in queue (e.g., #4 of 8)', G),
    ('Serving Number', 'Current patient being served (e.g., #3)', A),
    ('Total in Queue', 'Total patients waiting for that doctor/slot', T),
    ('Estimated Wait', 'AI-calculated wait time based on avg consultation', P),
    ('Live Updates', 'Real-time position changes via WebSocket push', G),
    ('Push Notifications', 'Alerts when it\'s almost your turn', B),
    ('Appointment Info', 'Doctor name, time slot, date, clinic address', GY),
    ('Digital Receipt', 'Payment receipt and summary after visit', R),
]

y = Inches(2.2)
for name, desc, color in app_features:
    add_rounded_rect(slide, Inches(6.8), y, Inches(6), Inches(0.5), W, RGBColor(230, 230, 235))
    add_rect(slide, Inches(6.8), y, Inches(0.06), Inches(0.5), color)
    add_circle(slide, Inches(7.0), y + Inches(0.08), Inches(0.3), color)
    add_text(slide, Inches(7.0), y + Inches(0.1), Inches(0.3), Inches(0.25), name[0], 10, W, True, PP_ALIGN.CENTER)
    add_text(slide, Inches(7.4), y + Inches(0.05), Inches(2), Inches(0.2), name, 9, D, True)
    add_text(slide, Inches(7.4), y + Inches(0.25), Inches(5), Inches(0.2), desc, 8, GY)
    y += Inches(0.55)

# Sync indicator
add_rounded_rect(slide, Inches(0.4), Inches(6.7), Inches(12.5), Inches(0.5), RGBColor(255, 248, 230), A)
add_text(slide, Inches(0.6), Inches(6.72), Inches(12), Inches(0.45),
    'SYNC: Web App status changes push to Android App via Socket.io in <100ms  |  Patient sees live queue position updates instantly',
    9, A, True, PP_ALIGN.LEFT)

# ================================================================
# SLIDE 4: TAM / SAM / SOM (DETAILED)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, Inches(0), Inches(0), SW, Inches(0), W)
topbar(slide, 'TAM / SAM / SOM', 'Market sizing for CatchQ - Total Addressable, Serviceable, and Obtainable Market')
botbar(slide, 4, 20)

# TAM - Total Addressable Market
add_rounded_rect(slide, Inches(0.4), Inches(1.6), Inches(4), Inches(5.2), LB, B)
add_rect(slide, Inches(0.4), Inches(1.6), Inches(4), Inches(0.6), B)
add_text(slide, Inches(0.6), Inches(1.65), Inches(3.5), Inches(0.5), 'TAM - Total Addressable Market', 13, W, True)
add_text(slide, Inches(0.6), Inches(2.3), Inches(3.5), Inches(0.5), 'Rs. 50,000 Cr', 28, B, True, PP_ALIGN.CENTER)
add_text(slide, Inches(0.6), Inches(2.8), Inches(3.5), Inches(0.3), 'Total clinic IT spend in India', 10, GY, align=PP_ALIGN.CENTER)
add_rect(slide, Inches(0.6), Inches(3.2), Inches(3.5), Inches(0.03), B)

add_multiline(slide, Inches(0.6), Inches(3.4), Inches(3.3), Inches(3.2), [
    'WHO:',
    '10 Lakh+ clinics across India',
    '30,000+ hospitals',
    '50 Crore+ patients visiting annually',
    '',
    'WHAT:',
    'Queue management software',
    'Patient engagement platform',
    'Billing & scheduling tools',
    'Bed management systems',
    '',
    'WHY:',
    '90% clinics use paper/WhatsApp',
    'Zero digital infrastructure',
    'Massive inefficiency = massive opportunity',
], 8, D)

# SAM - Serviceable Addressable Market
add_rounded_rect(slide, Inches(4.7), Inches(1.6), Inches(4), Inches(5.2), LG, G)
add_rect(slide, Inches(4.7), Inches(1.6), Inches(4), Inches(0.6), G)
add_text(slide, Inches(4.9), Inches(1.65), Inches(3.5), Inches(0.5), 'SAM - Serviceable Addressable', 13, W, True)
add_text(slide, Inches(4.9), Inches(2.3), Inches(3.5), Inches(0.5), 'Rs. 5,000 Cr', 28, G, True, PP_ALIGN.CENTER)
add_text(slide, Inches(4.9), Inches(2.8), Inches(3.5), Inches(0.3), 'SaaS-able segment', 10, GY, align=PP_ALIGN.CENTER)
add_rect(slide, Inches(4.9), Inches(3.2), Inches(3.5), Inches(0.03), G)

add_multiline(slide, Inches(4.9), Inches(3.4), Inches(3.3), Inches(3.2), [
    'SEGMENT:',
    'Clinics with 1-10 doctors',
    'Urban & semi-urban areas',
    'Tech-ready practices',
    '',
    'PRICING:',
    'Starter: Rs. 2,999/month',
    'Pro: Rs. 5,999/month',
    'Enterprise: Rs. 9,999/month',
    '',
    'CHANNELS:',
    'WhatsApp-first onboarding',
    'Free trial (14 days)',
    'Referral program',
], 8, D)

# SOM - Serviceable Obtainable Market
add_rounded_rect(slide, Inches(9), Inches(1.6), Inches(4), Inches(5.2), RGBColor(255, 248, 230), A)
add_rect(slide, Inches(9), Inches(1.6), Inches(4), Inches(0.6), A)
add_text(slide, Inches(9.2), Inches(1.65), Inches(3.5), Inches(0.5), 'SOM - Serviceable Obtainable', 13, W, True)
add_text(slide, Inches(9.2), Inches(2.3), Inches(3.5), Inches(0.5), 'Rs. 500 Cr', 28, A, True, PP_ALIGN.CENTER)
add_text(slide, Inches(9.2), Inches(2.8), Inches(3.5), Inches(0.3), 'Year 3 target', 10, GY, align=PP_ALIGN.CENTER)
add_rect(slide, Inches(9.2), Inches(3.2), Inches(3.5), Inches(0.03), A)

add_multiline(slide, Inches(9.2), Inches(3.4), Inches(3.3), Inches(3.2), [
    'YEAR 1:',
    '50 clinics  |  Rs. 36L ARR',
    '3 cities  |  5 team members',
    '',
    'YEAR 2:',
    '200 clinics  |  Rs. 1.8 Cr ARR',
    '10 cities  |  15 team members',
    '',
    'YEAR 3:',
    '1,000 clinics  |  Rs. 9.6 Cr ARR',
    '25 cities  |  40 team members',
    '',
    'GROWTH: 60% YoY',
    'RETENTION: 85%',
    'LTV/CAC: 5x',
], 8, D)

# ================================================================
# SLIDE 5: TAM SAM SOM - VISUAL
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, Inches(0), Inches(0), SW, Inches(0), W)
topbar(slide, 'Market Sizing Visual', 'Concentric circles showing TAM -> SAM -> SOM for CatchQ')
botbar(slide, 5, 20)

# Concentric circles
cx = Inches(6.666)
cy = Inches(4)

# TAM circle (outermost)
add_circle(slide, cx - Inches(2.8), cy - Inches(2.8), Inches(5.6), LB)
add_text(slide, cx - Inches(2.8), cy - Inches(2.6), Inches(5.6), Inches(0.4), 'TAM: Rs. 50,000 Cr', 14, B, True, PP_ALIGN.CENTER)
add_text(slide, cx - Inches(2.8), cy - Inches(2.2), Inches(5.6), Inches(0.3), 'All healthcare IT in India', 9, GY, align=PP_ALIGN.CENTER)

# SAM circle (middle)
add_circle(slide, cx - Inches(2), cy - Inches(2), Inches(4), LG)
add_text(slide, cx - Inches(2), cy - Inches(1.8), Inches(4), Inches(0.4), 'SAM: Rs. 5,000 Cr', 14, G, True, PP_ALIGN.CENTER)
add_text(slide, cx - Inches(2), cy - Inches(1.4), Inches(4), Inches(0.3), 'SaaS clinics (1-10 doctors)', 9, GY, align=PP_ALIGN.CENTER)

# SOM circle (innermost)
add_circle(slide, cx - Inches(1.2), cy - Inches(1.2), Inches(2.4), RGBColor(255, 248, 230))
add_text(slide, cx - Inches(1.2), cy - Inches(1.0), Inches(2.4), Inches(0.4), 'SOM: Rs. 500 Cr', 14, A, True, PP_ALIGN.CENTER)
add_text(slide, cx - Inches(1.2), cy - Inches(0.6), Inches(2.4), Inches(0.3), 'Year 3 target', 9, GY, align=PP_ALIGN.CENTER)

# Left: market breakdown
add_text(slide, Inches(0.5), Inches(1.6), Inches(3), Inches(0.4), 'MARKET BREAKDOWN', 12, B, True)
add_rect(slide, Inches(0.5), Inches(1.95), Inches(1.5), Inches(0.04), G)

breakdown = [
    ('Clinics', '10 Lakh+', '90% digital-free', B),
    ('Hospitals', '30,000+', 'Need bed management', G),
    ('Patients', '50 Crore+', 'Visit clinics annually', A),
    ('Queue Mgmt', 'Rs. 1.2B', 'Global market 2028', T),
    ('Clinic Mgmt', 'Rs. 8.4B', 'Global market 2028', P),
    ('AI Healthcare', 'Rs. 45B', 'Global market 2030', R),
]

y = Inches(2.2)
for name, val, desc, color in breakdown:
    add_rounded_rect(slide, Inches(0.5), y, Inches(3.5), Inches(0.5), W, RGBColor(230, 230, 235))
    add_rect(slide, Inches(0.5), y, Inches(0.06), Inches(0.5), color)
    add_text(slide, Inches(0.7), y + Inches(0.05), Inches(1.5), Inches(0.2), name, 9, D, True)
    add_text(slide, Inches(0.7), y + Inches(0.25), Inches(1.5), Inches(0.2), val, 9, color, True)
    add_text(slide, Inches(2.3), y + Inches(0.1), Inches(1.5), Inches(0.3), desc, 7.5, GY)
    y += Inches(0.55)

# Right: competitive positioning
add_text(slide, Inches(10.5), Inches(1.6), Inches(3), Inches(0.4), 'WHY CATCHQ WINS', 12, B, True)
add_rect(slide, Inches(10.5), Inches(1.95), Inches(1.5), Inches(0.04), G)

advantages = [
    'Only all-in-one platform',
    'WhatsApp-first (India native)',
    'AI from Day 1 (not bolt-on)',
    'Rs. 2,999/month (affordable)',
    'Real-time queue (Socket.io)',
    'Bed management built-in',
    '30+ APIs already built',
    '2,100+ clinics in database',
]

y = Inches(2.2)
for adv in advantages:
    add_rounded_rect(slide, Inches(10.5), y, Inches(2.5), Inches(0.35), LG)
    add_rect(slide, Inches(10.5), y, Inches(0.06), Inches(0.35), G)
    add_text(slide, Inches(10.7), y + Inches(0.05), Inches(2.2), Inches(0.25), adv, 8, D, True)
    y += Inches(0.4)

# Bottom metrics
add_rounded_rect(slide, Inches(0.4), Inches(6.7), Inches(12.5), Inches(0.5), LB, B)
add_text(slide, Inches(0.6), Inches(6.72), Inches(12), Inches(0.45),
    'CATCHQ: India\'s first WhatsApp-native, AI-powered, all-in-one clinic operating system  |  10 Lakh+ addressable clinics  |  Rs. 50,000 Cr TAM',
    9, B, True, PP_ALIGN.LEFT)

# ================================================================
# SAVE
# ================================================================
prs.save(OUT)
print(f"PPTX generated: {OUT}")
print("5 slides: Queue Overview, Status Flow, Web vs App, TAM/SAM/SOM, Market Visual")
