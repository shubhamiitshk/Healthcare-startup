from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# Brand colors
G = RGBColor(27, 186, 141)
B = RGBColor(22, 71, 114)
W = RGBColor(255, 255, 255)
D = RGBColor(40, 40, 50)
GY = RGBColor(130, 130, 140)
LB = RGBColor(230, 242, 255)
LG = RGBColor(240, 248, 240)
LGR = RGBColor(230, 250, 230)
R = RGBColor(220, 60, 60)
P = RGBColor(140, 60, 180)
A = RGBColor(230, 160, 30)
T = RGBColor(60, 130, 180)

IMG_DIR = r"C:\Users\sumit\Downloads\Catchq\images"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SW = Inches(13.333)
SH = Inches(7.5)

def add_bg(slide, color=W):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_circle(slide, x, y, size, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text(slide, x, y, w, h, text, size=12, color=D, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    p.font.name = 'Calibri'
    return txBox

def add_multiline(slide, x, y, w, h, lines, size=10, color=D, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = align
        p.font.name = 'Calibri'
        p.space_after = Pt(4)
    return txBox

def add_image_safe(slide, img_name, x, y, w, h):
    path = os.path.join(IMG_DIR, img_name)
    if os.path.exists(path):
        try:
            pic = slide.shapes.add_picture(path, x, y, w, h)
            # Crop to fit
            return pic
        except:
            pass
    return None

def add_image_overlay(slide, img_name, x, y, w, h, overlay_alpha=0.55):
    """Add image with dark overlay for text readability"""
    path = os.path.join(IMG_DIR, img_name)
    if os.path.exists(path):
        try:
            pic = slide.shapes.add_picture(path, x, y, w, h)
            overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = RGBColor(22, 51, 90)
            overlay.line.fill.background()
            return pic
        except:
            pass
    return None

def add_image_card(slide, img_name, x, y, w, h, title, desc, color):
    """Card with image on top and text below"""
    add_rounded_rect(slide, x, y, w, h, W, RGBColor(230, 230, 235))
    pic = add_image_safe(slide, img_name, x + Inches(0.1), y + Inches(0.1), w - Inches(0.2), h * 0.55)
    if not pic:
        add_rect(slide, x + Inches(0.1), y + Inches(0.1), w - Inches(0.2), h * 0.55, RGBColor(240, 240, 245))
    add_rect(slide, x, y + h * 0.55, w, Inches(0.05), color)
    add_text(slide, x + Inches(0.15), y + h * 0.58, w - Inches(0.3), Inches(0.35), title, 11, color, True)
    add_text(slide, x + Inches(0.15), y + h * 0.58 + Inches(0.35), w - Inches(0.3), Inches(0.8), desc, 8.5, GY)

def topbar(slide, title, subtitle):
    add_rect(slide, Inches(0), Inches(0), SW, Inches(1.2), B)
    add_rect(slide, Inches(0), Inches(1.2), SW, Inches(0.08), G)
    add_text(slide, Inches(0.5), Inches(0.15), Inches(10), Inches(0.7), title, 28, W, True)
    add_text(slide, Inches(0.5), Inches(0.75), Inches(10), Inches(0.4), subtitle, 13, RGBColor(170, 200, 230))

def botbar(slide, num, total=18):
    add_rect(slide, Inches(0), SH - Inches(0.35), SW, Inches(0.35), RGBColor(20, 55, 90))
    add_text(slide, Inches(0.5), SH - Inches(0.33), Inches(6), Inches(0.3),
        'CatchQ  |  IIT Kharagpur Platinum Jubilee Seed Fund Application', 8, RGBColor(140, 160, 180))
    add_text(slide, SW - Inches(1.5), SH - Inches(0.33), Inches(1.2), Inches(0.3),
        f'{num} / {total}', 8, RGBColor(140, 160, 180), align=PP_ALIGN.RIGHT)

def stat_box(slide, x, y, w, h, val, label, color):
    add_rounded_rect(slide, x, y, w, h, W, RGBColor(230, 230, 235))
    add_rect(slide, x, y, w, Inches(0.06), color)
    add_text(slide, x, y + Inches(0.15), w, Inches(0.45), val, 24, color, True, PP_ALIGN.CENTER)
    add_multiline(slide, x + Inches(0.1), y + Inches(0.6), w - Inches(0.2), Inches(0.5),
        label.split('\n'), 9, GY, align=PP_ALIGN.CENTER)

def image_card(slide, img_name, x, y, w, h, title, desc, color):
    """Card with image on top and text below"""
    # card background
    add_rounded_rect(slide, x, y, w, h, W, RGBColor(230, 230, 235))
    # image area
    pic = add_image_safe(slide, img_name, x + Inches(0.1), y + Inches(0.1), w - Inches(0.2), h * 0.55)
    if not pic:
        add_rect(slide, x + Inches(0.1), y + Inches(0.1), w - Inches(0.2), h * 0.55, RGBColor(240, 240, 245))
        add_text(slide, x + Inches(0.1), y + Inches(0.3), w - Inches(0.2), Inches(0.4),
            '[Image]', 10, GY, align=PP_ALIGN.CENTER)
    # color accent
    add_rect(slide, x, y + h * 0.55, w, Inches(0.05), color)
    # title
    add_text(slide, x + Inches(0.15), y + h * 0.58, w - Inches(0.3), Inches(0.35), title, 11, color, True)
    # desc
    add_text(slide, x + Inches(0.15), y + h * 0.58 + Inches(0.35), w - Inches(0.3), Inches(0.8), desc, 8.5, GY)

def process_bar(slide, x, y, w, color, label, desc):
    """Horizontal process bar with icon"""
    add_rounded_rect(slide, x, y, w, Inches(0.6), RGBColor(248, 248, 252))
    add_rect(slide, x, y, Inches(0.08), Inches(0.6), color)
    add_circle(slide, x + Inches(0.2), y + Inches(0.1), Inches(0.4), color)
    add_text(slide, x + Inches(0.2), y + Inches(0.12), Inches(0.4), Inches(0.35), label[0], 12, W, True, PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.7), y + Inches(0.05), w - Inches(0.9), Inches(0.25), label, 9, D, True)
    add_text(slide, x + Inches(0.7), y + Inches(0.3), w - Inches(0.9), Inches(0.25), desc, 7.5, GY)

# ================================================================
# SLIDE 1: TITLE (with hero image)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, B)

# Background image with overlay
add_image_overlay(slide, 'digital_health.jpg', Inches(0), Inches(0), SW, SH, 0.65)

# decorative circles
add_circle(slide, Inches(-1), Inches(-1), Inches(4), RGBColor(30, 85, 140))
add_circle(slide, Inches(10), Inches(4.5), Inches(3.5), RGBColor(20, 60, 110))

add_rect(slide, Inches(0), SH - Inches(1), SW, Inches(1), G)
add_rect(slide, Inches(0), SH - Inches(1.08), SW, Inches(0.08), RGBColor(20, 140, 100))

add_text(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.6), 'CATCHQ', 48, W, True)
add_text(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.8), "India's Clinic & Hospital Operating System", 28, G, True)
add_rect(slide, Inches(0.8), Inches(2.3), Inches(2.5), Inches(0.06), G)
add_text(slide, Inches(0.8), Inches(2.7), Inches(6), Inches(1.2),
    "WhatsApp booking, AI-powered diagnosis, real-time queue management,\ndigital billing, pharmacy integration, and hospital bed tracking\n- all in one platform for India's 10 Lakh+ clinics.",
    15, RGBColor(200, 220, 240))
add_text(slide, Inches(0.8), Inches(4.2), Inches(10), Inches(0.5), 'Platinum Jubilee Seed Fund Application', 20, W, True)
add_text(slide, Inches(0.8), Inches(4.8), Inches(5), Inches(0.4), '[Your Name]  |  [Your Email]  |  [Your Phone]', 12, RGBColor(180, 200, 220))
add_text(slide, Inches(0.8), Inches(5.2), Inches(5), Inches(0.4), 'IIT Kharagpur  |  [Your Department]', 12, RGBColor(180, 200, 220))

# Right side: floating phone mockup with image
add_rounded_rect(slide, Inches(9), Inches(1.5), Inches(3.5), Inches(5.5), RGBColor(30, 30, 35))
add_rect(slide, Inches(9.15), Inches(1.8), Inches(3.2), Inches(4.8), W)
add_image_safe(slide, 'doctor_illustration.png', Inches(9.4), Inches(2.0), Inches(2.6), Inches(2.0))
add_rect(slide, Inches(9.15), Inches(4.2), Inches(3.2), Inches(0.3), G)
add_text(slide, Inches(9.15), Inches(4.25), Inches(3.2), Inches(0.25), 'CATCHQ', 10, W, True, PP_ALIGN.CENTER)
add_text(slide, Inches(9.3), Inches(4.6), Inches(2.9), Inches(1.8),
    "Book Appointment\nAI Doctor Match\nReal-time Queue\nDigital Receipt\nLab Reports", 9, D, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0.8), SH - Inches(0.85), Inches(4), Inches(0.3), 'CATCHQ', 16, W, True)
add_text(slide, Inches(0.8), SH - Inches(0.55), Inches(8), Inches(0.3),
    "Healthcare shouldn't start with waiting. It should start with CatchQ.", 11, RGBColor(200, 240, 220))

# ================================================================
# SLIDE 2: PROBLEM (with background image)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, W)
topbar(slide, 'The Problem', 'India manages healthcare with WhatsApp messages, phone calls, and paper registers')
botbar(slide, 2, 18)

# Left side: hero image
add_image_overlay(slide, 'hospital_waiting.jpg', Inches(0.4), Inches(1.6), Inches(5.5), Inches(3.5), 0.5)
add_text(slide, Inches(0.6), Inches(2.5), Inches(5), Inches(0.5), 'THE CRISIS', 22, W, True, PP_ALIGN.CENTER)
add_text(slide, Inches(0.6), Inches(3.1), Inches(5), Inches(0.8),
    "India's 10 Lakh+ clinics operate\nwith broken, manual systems", 14, RGBColor(220, 240, 255), align=PP_ALIGN.CENTER)

problems = [
    ('Missed Calls', "30-40% bookings lost\nto missed calls", R),
    ('No Digital Booking', "90% clinics have no\nonline booking system", A),
    ('Zero Visibility', "Patients wait hours\nwith no queue info", P),
    ('No-Show Crisis', "40% no-show rate.\nRs. 5L+ monthly loss", R),
    ('Paper Records', "No digital trail.\nZero analytics.", GY),
    ('Billing Chaos', "Manual billing.\nCash-only systems.", RGBColor(180, 140, 60)),
    ('Staff Overload', "60% time on phone\ncalls and confirmations", T),
    ('Zero Intelligence', "No patient analytics.\nNo follow-ups.", RGBColor(100, 140, 180)),
]

for i, (title, desc, color) in enumerate(problems):
    col = i % 4
    row = i // 4
    x = Inches(6.3 + col * 1.75)
    y = Inches(1.6 + row * 1.8)

    add_rounded_rect(slide, x, y, Inches(1.6), Inches(1.5), RGBColor(250, 250, 252))
    add_rect(slide, x, y, Inches(0.06), Inches(1.5), color)
    add_circle(slide, x + Inches(0.1), y + Inches(0.1), Inches(0.3), color)
    add_text(slide, x + Inches(0.1), y + Inches(0.12), Inches(0.3), Inches(0.25), str(i+1), 10, W, True, PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.45), y + Inches(0.1), Inches(1.05), Inches(0.3), title, 8, color, True)
    add_text(slide, x + Inches(0.1), y + Inches(0.55), Inches(1.4), Inches(0.85), desc, 7.5, GY)

add_rounded_rect(slide, Inches(0.4), Inches(5.8), Inches(5.5), Inches(0.5), RGBColor(255, 245, 245), R)
add_text(slide, Inches(0.6), Inches(5.82), Inches(5), Inches(0.45),
    "Impact: Rs. 5-10 Lakh lost per doctor per year to operational inefficiency",
    9, R, True, PP_ALIGN.LEFT)

# ================================================================
# SLIDE 3: SOLUTION (with images)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, W)
topbar(slide, 'Our Solution', 'CatchQ replaces 7 disconnected tools with one unified platform')
botbar(slide, 3, 18)

# Left image: clinic
add_image_card(slide, 'hospital_illustration.png', Inches(0.4), Inches(1.6), Inches(4), Inches(3.2),
    'CLINIC SIDE', 'WhatsApp Bot, Dashboard, Queue Board,\nSmart TV, Billing, Scheduling, Patient DB', B)

# Right image: patient
add_image_card(slide, 'nurse_illustration.png', Inches(4.7), Inches(1.6), Inches(4), Inches(3.2),
    'PATIENT SIDE', 'WhatsApp Booking, Mobile App, Queue Tracking,\nDigital Receipt, Lab AI, Follow-up Reminders', G)

# Center: tech stack
add_rounded_rect(slide, Inches(9), Inches(1.6), Inches(4), Inches(3.2), RGBColor(245, 248, 255))
add_text(slide, Inches(9.2), Inches(1.7), Inches(3.5), Inches(0.35), 'TECHNOLOGY STACK', 11, B, True)
add_rect(slide, Inches(9.2), Inches(2.0), Inches(1.2), Inches(0.03), B)

tech_list = [
    ('Next.js 15', 'Frontend', B), ('NestJS', 'Backend', B),
    ('React Native', 'Mobile', G), ('Firebase', 'Auth', G),
    ('Socket.io', 'Real-time', A), ('PostgreSQL', 'Database', P),
    ('Tailwind', 'UI', T), ('TypeScript', 'Language', B),
    ('Python AI', 'ML/AI', G), ('WhatsApp API', 'Chat', G),
]
ty = Inches(2.15)
for name, sub, color in tech_list:
    add_circle(slide, Inches(9.3), ty, Inches(0.25), color)
    add_text(slide, Inches(9.3), ty + Inches(0.02), Inches(0.25), Inches(0.2), name[0], 8, W, True, PP_ALIGN.CENTER)
    add_text(slide, Inches(9.65), ty, Inches(1.5), Inches(0.2), name, 8, D, True)
    add_text(slide, Inches(11.2), ty, Inches(1.5), Inches(0.2), sub, 8, GY)
    ty += Inches(0.3)

# Bottom: solution summary with image
add_rect(slide, Inches(0.4), Inches(5.1), SW - Inches(0.8), Inches(0.03), RGBColor(230, 230, 235))
add_image_safe(slide, 'whatsapp_booking.jpg', Inches(0.4), Inches(5.3), Inches(4), Inches(1.5))
add_rounded_rect(slide, Inches(4.7), Inches(5.3), Inches(8.5), Inches(1.5), LG, G)
add_text(slide, Inches(4.9), Inches(5.4), Inches(8), Inches(0.35),
    "7 core problems, 1 unified solution", 14, G, True)
add_text(slide, Inches(4.9), Inches(5.8), Inches(8), Inches(0.9),
    "Every feature is already built and operational. From WhatsApp booking to AI-powered\ndiagnostics, from real-time queue management to hospital bed tracking - CatchQ is\nthe complete operating system for Indian healthcare.", 10, D)

# ================================================================
# SLIDE 4: ARCHITECTURE
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, W)
topbar(slide, 'Platform Architecture', 'Full-stack system built for scale, security, and speed')
botbar(slide, 4, 18)

# 3-tier architecture
tiers = [
    ('PATIENT TOUCHPOINTS', [
        ('WhatsApp Bot', 'Book via chat', G),
        ('Mobile App', '11 screens', B),
        ('Web Booking', 'Landing page', A),
        ('Walk-in', 'Clinic kiosk', P),
    ], LB, Inches(0.4)),
    ('BACKEND CORE', [
        ('NestJS API', '6 modules', B),
        ('Firebase Auth', 'JWT + OTP', G),
        ('Socket.io', 'Real-time', A),
        ('PostgreSQL', '7 tables', P),
    ], RGBColor(245, 248, 255), Inches(4.6)),
    ('FEATURES', [
        ('Queue Management', 'Live updates', G),
        ('Billing Engine', 'Auto-invoice', R),
        ('AI Analytics', 'Smart insights', A),
        ('Bed Tracking', 'Real-time', T),
    ], RGBColor(255, 248, 245), Inches(9)),
]

for title, items, bg_color, x in tiers:
    add_rounded_rect(slide, x, Inches(1.6), Inches(3.8), Inches(2.8), bg_color)
    add_text(slide, x + Inches(0.2), Inches(1.7), Inches(3.4), Inches(0.4), title, 11, B, True)
    add_rect(slide, x + Inches(0.2), Inches(2.05), Inches(1.2), Inches(0.03), B)
    ty = Inches(2.2)
    for name, sub, color in items:
        add_circle(slide, x + Inches(0.3), ty, Inches(0.35), color)
        add_text(slide, x + Inches(0.3), ty + Inches(0.03), Inches(0.35), Inches(0.3), name[0], 10, W, True, PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.75), ty, Inches(2.5), Inches(0.25), name, 10, D, True)
        add_text(slide, x + Inches(0.75), ty + Inches(0.22), Inches(2.5), Inches(0.2), sub, 8, GY)
        ty += Inches(0.45)

# Arrows between tiers
for ax in [Inches(4.2), Inches(8.4)]:
    add_text(slide, ax, Inches(2.6), Inches(0.4), Inches(0.4), '>>', 14, GY, True, PP_ALIGN.CENTER)

# Platform scale stats
add_text(slide, Inches(0.5), Inches(4.7), Inches(3), Inches(0.35), 'PLATFORM SCALE', 11, B, True)
add_rect(slide, Inches(0.5), Inches(5.0), Inches(1.2), Inches(0.03), G)

scale_items = [
    ('30+', 'API Endpoints', B), ('14', 'Frontend Pages', G), ('11', 'Mobile Screens', A),
    ('7', 'Database Tables', P), ('22', 'Landing Parts', T), ('6', 'Backend Modules', R),
]
sx = Inches(0.5)
for val, lbl, color in scale_items:
    w = Inches(2)
    add_rounded_rect(slide, sx, Inches(5.2), w, Inches(0.65), RGBColor(248, 248, 252))
    add_text(slide, sx, Inches(5.22), w, Inches(0.35), val, 18, color, True, PP_ALIGN.CENTER)
    add_text(slide, sx, Inches(5.55), w, Inches(0.3), lbl, 9, GY, align=PP_ALIGN.CENTER)
    sx += w + Inches(0.1)

# ================================================================
# SLIDE 5: HOW IT WORKS (with phone mockups + images)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, W)
topbar(slide, 'How It Works', '8-step patient journey from booking to pharmacy refill')
botbar(slide, 5, 18)

steps = [
    ('01', 'BOOK', 'WhatsApp or app.\nAI suggests doctor.', G, 'whatsapp_icon.png'),
    ('02', 'CONFIRM', 'Instant confirmation.\nCalendar sync.', B, 'appointment_icon.png'),
    ('03', 'TRACK', 'Real-time queue.\n"You are #3."', RGBColor(27, 120, 90), 'queue_icon.png'),
    ('04', 'VISIT', 'Smart display.\nLive patient list.', A, 'doctor_illustration.png'),
    ('05', 'PAY', 'Digital receipt.\nUPI/Card/Cash.', P, 'medical_history.png'),
    ('06', 'REPORT', 'AI reads lab report.\nSimple language.', R, 'ai_brain.png'),
    ('07', 'FOLLOW', 'Auto reminder.\nNext appointment.', T, 'health_checkup.png'),
    ('08', 'REFILL', 'Pharmacy alert.\nFamily notified.', RGBColor(100, 60, 140), 'pharmacy_icon.png'),
]

for i, (num, title, desc, color, img) in enumerate(steps):
    col = i % 4
    row = i // 4
    x = Inches(0.4 + col * 3.2)
    y = Inches(1.6 + row * 2.8)

    # phone frame
    add_rounded_rect(slide, x, y, Inches(2.9), Inches(2.5), RGBColor(45, 45, 50))
    # screen with image
    add_rect(slide, x + Inches(0.1), y + Inches(0.1), Inches(2.7), Inches(1.2), W)
    add_image_safe(slide, img, x + Inches(0.1), y + Inches(0.1), Inches(2.7), Inches(1.2))
    # dark overlay on image
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.1), y + Inches(0.1), Inches(2.7), Inches(1.2))
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(22, 40, 70)
    # overlay (solid color)
    overlay.line.fill.background()
    # step circle on image
    add_circle(slide, x + Inches(1.0), y + Inches(0.35), Inches(0.6), color)
    add_text(slide, x + Inches(1.0), y + Inches(0.4), Inches(0.6), Inches(0.5), num, 16, W, True, PP_ALIGN.CENTER)
    # app content area
    add_rect(slide, x + Inches(0.1), y + Inches(1.3), Inches(2.7), Inches(0.2), color)
    add_text(slide, x + Inches(0.1), y + Inches(1.32), Inches(2.7), Inches(0.18), 'CATCHQ', 7, W, True, PP_ALIGN.CENTER)
    # title and desc
    add_text(slide, x + Inches(0.1), y + Inches(1.6), Inches(2.7), Inches(0.3), title, 12, color, True, PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.2), y + Inches(1.95), Inches(2.5), Inches(0.5), desc, 8.5, GY, align=PP_ALIGN.CENTER)

# Bottom tagline
add_rect(slide, Inches(0), SH - Inches(0.8), SW, Inches(0.04), LG)
add_text(slide, Inches(0), SH - Inches(0.7), SW, Inches(0.35), 'PATIENT JOURNEY', 14, B, True, PP_ALIGN.CENTER)
add_text(slide, Inches(0), SH - Inches(0.4), SW, Inches(0.3),
    'Every touchpoint is digital, real-time, and designed for zero friction', 11, GY, align=PP_ALIGN.CENTER)

# ================================================================
# SLIDE 6: DASHBOARD FEATURES (with screenshot)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, W)
topbar(slide, 'Dashboard Features', 'Complete clinic management system - 8 modules, already built and operational')
botbar(slide, 6, 18)

# Left: Dashboard screenshot area
add_rounded_rect(slide, Inches(0.4), Inches(1.6), Inches(5.5), Inches(5.2), RGBColor(245, 248, 255))
add_text(slide, Inches(0.6), Inches(1.7), Inches(5), Inches(0.35), 'ADMIN DASHBOARD', 12, B, True)
add_image_safe(slide, 'digital_health.jpg', Inches(0.6), Inches(2.1), Inches(5.1), Inches(2.5))
# Dashboard feature boxes
dash_features = [
    ('Real-time KPIs', 'Patients, Queue, Revenue', G),
    ('Doctor Management', 'Schedule, Availability', B),
    ('Patient Database', 'Search, History, Records', A),
]
dy = Inches(4.8)
for name, desc, color in dash_features:
    add_rounded_rect(slide, Inches(0.6), dy, Inches(5.1), Inches(0.55), W, RGBColor(230, 230, 235))
    add_rect(slide, Inches(0.6), dy, Inches(0.06), Inches(0.55), color)
    add_circle(slide, Inches(0.8), dy + Inches(0.08), Inches(0.35), color)
    add_text(slide, Inches(0.8), dy + Inches(0.1), Inches(0.35), Inches(0.3), name[0], 10, W, True, PP_ALIGN.CENTER)
    add_text(slide, Inches(1.25), dy + Inches(0.05), Inches(2), Inches(0.25), name, 9, D, True)
    add_text(slide, Inches(1.25), dy + Inches(0.28), Inches(4), Inches(0.25), desc, 8, GY)
    dy += Inches(0.6)

# Right: feature cards
features = [
    ('QUEUE', 'Per-doctor live queue.\nSocket.io updates.', G, 'Real-time'),
    ('BILLING', 'Invoice, receipts.\nCash/UPI/Card.', R, '3 Modes'),
    ('FOLLOW-UPS', 'Cron reminders.\nRetention workflows.', P, 'Automated'),
    ('PATIENT APP', 'OTP login. Queue.\nMobile booking.', T, '11 Screens'),
    ('LANDING PAGE', '22 components.\nPricing, WhatsApp.', GY, '22 Parts'),
]

fy = Inches(1.6)
for title, desc, color, badge in features:
    add_rounded_rect(slide, Inches(6.2), fy, Inches(6.8), Inches(0.95), W, RGBColor(230, 230, 235))
    add_rect(slide, Inches(6.2), fy, Inches(0.06), Inches(0.95), color)
    add_circle(slide, Inches(6.4), fy + Inches(0.15), Inches(0.5), color)
    add_text(slide, Inches(6.4), fy + Inches(0.18), Inches(0.5), Inches(0.45), title[0], 12, W, True, PP_ALIGN.CENTER)
    add_text(slide, Inches(7.0), fy + Inches(0.1), Inches(2), Inches(0.3), title, 11, D, True)
    add_text(slide, Inches(7.0), fy + Inches(0.4), Inches(4), Inches(0.5), desc, 8.5, GY)
    # badge
    bw = Inches(len(badge) * 0.1 + 0.3)
    add_rounded_rect(slide, Inches(12.5) - bw, fy + Inches(0.1), bw, Inches(0.3), color)
    add_text(slide, Inches(12.5) - bw, fy + Inches(0.12), bw, Inches(0.26), badge, 8, W, True, PP_ALIGN.CENTER)
    fy += Inches(1.05)

# ================================================================
# SLIDE 7: AI INNOVATION (with AI imagery)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, W)
topbar(slide, 'AI Innovation', 'Four AI capabilities that will transform Indian healthcare SaaS')
botbar(slide, 7, 18)

ai_features = [
    ('SMART DOCTOR\nMATCHING', 'Patient describes symptoms\n-> AI recommends right\nspecialist immediately\n\nImpact: Reduces wrong\nconsultations by 40%', B, 'Symptom -> Specialist', 'doctor_illustration.png'),
    ('REPORT ANALYSIS\n(RAG)', 'Lab report uploaded\n-> explains in simple\nlanguage\n\nImpact: 90% patients\nunderstand results', G, 'Report -> Plain Language', 'medical_history.png'),
    ('VOICE HEALTH\nASSISTANT', 'Patient calls. AI answers,\ntriages, books appointment\nautomatically\n\nImpact: 24/7 availability.\nZero missed calls.', A, 'Call -> Triage -> Booking', 'whatsapp_icon.png'),
    ('PRESCRIPTION\nDIGITIZER', 'Handwritten prescription\n-> AI converts to digital\n-> sends to pharmacy\n\nImpact: 100% accuracy.\nZero errors.', P, 'Handwriting -> Pharmacy', 'digital_health_icon.png'),
]

for i, (title, desc, color, badge, img) in enumerate(ai_features):
    x = Inches(0.4 + i * 3.2)
    y = Inches(1.6)

    # card
    add_rounded_rect(slide, x, y, Inches(3), Inches(5.2), RGBColor(248, 248, 252))
    # image area
    add_image_safe(slide, img, x + Inches(0.1), y + Inches(0.1), Inches(2.8), Inches(1.5))
    # overlay
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.1), y + Inches(0.1), Inches(2.8), Inches(1.5))
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = color
    # overlay (solid color)
    overlay.line.fill.background()
    # title on image
    add_text(slide, x + Inches(0.2), y + Inches(0.2), Inches(2.6), Inches(0.7), title, 13, W, True)
    # badge
    bw = Inches(len(badge) * 0.09 + 0.3)
    add_rounded_rect(slide, x + Inches(0.2), y + Inches(1.3), bw, Inches(0.25), W)
    add_text(slide, x + Inches(0.2), y + Inches(1.32), bw, Inches(0.22), badge, 8, color, True, PP_ALIGN.CENTER)
    # divider
    add_rect(slide, x + Inches(0.15), y + Inches(1.7), Inches(2.7), Inches(0.02), color)
    # desc
    add_text(slide, x + Inches(0.15), y + Inches(1.8), Inches(2.7), Inches(3.2), desc, 9.5, D)

# insight box
add_rounded_rect(slide, Inches(0.4), Inches(6.95), Inches(12.5), Inches(0.4), LGR, G)
add_text(slide, Inches(0.6), Inches(6.97), Inches(12), Inches(0.35),
    "Key Insight: AI in healthcare market will reach $45 Billion by 2030. CatchQ builds these capabilities from Day 1.",
    10, RGBColor(20, 100, 50), True, PP_ALIGN.LEFT)

# ================================================================
# SLIDE 8: HOSPITAL BED MANAGEMENT (with hospital image)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, W)
topbar(slide, 'Hospital Bed Management', 'Real-time bed tracking for emergency and ward operations')
botbar(slide, 8, 18)

# stat cards
stat_box(slide, Inches(0.4), Inches(1.6), Inches(2.9), Inches(0.9), '48', 'Total Beds', B)
stat_box(slide, Inches(3.6), Inches(1.6), Inches(2.9), Inches(0.9), '31', 'Occupied (65%)', R)
stat_box(slide, Inches(6.8), Inches(1.6), Inches(2.9), Inches(0.9), '12', 'Available (25%)', G)
stat_box(slide, Inches(10), Inches(1.6), Inches(2.9), Inches(0.9), '3', 'Emergency Free', A)

# Left: hospital image + features
add_image_safe(slide, 'bed_icon.png', Inches(1.5), Inches(2.9), Inches(3), Inches(1.8))
overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(2.7), Inches(5.5), Inches(2.2))
overlay.fill.solid()
overlay.fill.fore_color.rgb = RGBColor(22, 51, 90)
# overlay (solid color)
overlay.line.fill.background()
add_text(slide, Inches(0.6), Inches(3.2), Inches(5), Inches(0.4), 'REAL-TIME BED TRACKING', 14, W, True, PP_ALIGN.CENTER)
add_text(slide, Inches(0.6), Inches(3.7), Inches(5), Inches(0.8),
    "ICU, General, Emergency, Maternity, Pediatric\nWard-wise status with equipment tracking", 10, RGBColor(200, 220, 240), align=PP_ALIGN.CENTER)

# Right: bed grid mockup
add_text(slide, Inches(6.5), Inches(2.8), Inches(3), Inches(0.35), 'BED STATUS GRID', 11, B, True)
add_rect(slide, Inches(6.5), Inches(3.1), Inches(1.2), Inches(0.03), G)

wards_data = [
    ('ICU', 8, [1, 1, 0, 1, 0, 1, 1, 0]),
    ('General A', 10, [0, 1, 1, 0, 0, 1, 1, 0, 1, 0]),
    ('Emergency', 6, [1, 0, 1, 0, 1, 0]),
    ('Maternity', 6, [0, 0, 1, 1, 0, 0]),
    ('Pediatric', 6, [1, 0, 0, 1, 0, 1]),
]

y_pos = Inches(3.3)
for ward_name, count, statuses in wards_data:
    add_text(slide, Inches(6.5), y_pos, Inches(2), Inches(0.3), ward_name, 9, D, True)
    for j, occ in enumerate(statuses):
        bx = Inches(9) + Inches(j * 0.35)
        color = R if occ else G
        add_rect(slide, bx, y_pos, Inches(0.28), Inches(0.28), color)
    y_pos += Inches(0.45)

# Legend
add_text(slide, Inches(6.5), Inches(5.7), Inches(1), Inches(0.25), 'Legend:', 8, GY, True)
add_rect(slide, Inches(7.5), Inches(5.73), Inches(0.2), Inches(0.18), G)
add_text(slide, Inches(7.8), Inches(5.7), Inches(1), Inches(0.25), 'Available', 7, GY)
add_rect(slide, Inches(8.8), Inches(5.73), Inches(0.2), Inches(0.18), R)
add_text(slide, Inches(9.1), Inches(5.7), Inches(1), Inches(0.25), 'Occupied', 7, GY)
add_rect(slide, Inches(10.1), Inches(5.73), Inches(0.2), Inches(0.18), A)
add_text(slide, Inches(10.4), Inches(5.7), Inches(1), Inches(0.25), 'Reserved', 7, GY)

# Bottom features
bed_features = [
    ('Ward Grid', 'Color-coded status'),
    ('Emergency Queue', '< 5 min allocation'),
    ('Admission Track', 'Full lifecycle'),
    ('Family Alerts', 'Real-time notifications'),
    ('Equipment', 'Ventilator, monitor, O2'),
    ('Billing Link', 'Auto daily charges'),
]
fx = Inches(0.4)
for name, desc in bed_features:
    add_rounded_rect(slide, fx, Inches(6.0), Inches(2), Inches(0.7), LG)
    add_rect(slide, fx, Inches(6.0), Inches(2), Inches(0.06), G)
    add_text(slide, fx + Inches(0.1), Inches(6.1), Inches(1.8), Inches(0.25), name, 8, D, True)
    add_text(slide, fx + Inches(0.1), Inches(6.35), Inches(1.8), Inches(0.25), desc, 7, GY)
    fx += Inches(2.1)

# ================================================================
# SLIDE 9: TRACTION & METRICS (with images)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, W)
topbar(slide, 'Traction & Metrics', 'What we have built and the impact it delivers')
botbar(slide, 9, 18)

# Left: What's Built
add_text(slide, Inches(0.5), Inches(1.6), Inches(3), Inches(0.4), 'WHAT WE HAVE BUILT', 13, B, True)
add_rect(slide, Inches(0.5), Inches(1.95), Inches(1.5), Inches(0.04), G)

built_items = [
    ('2,100+', 'Clinics registered in database'),
    ('35+', 'Doctors across 15+ specialties'),
    ('180+', 'Weekly schedule slots'),
    ('11', 'Mobile app screens'),
    ('8', 'Admin dashboard pages'),
    ('30+', 'REST API endpoints'),
    ('Socket.io', 'Real-time WebSocket gateway'),
    ('22', 'Landing page components'),
]

y_pos = Inches(2.2)
for val, desc in built_items:
    add_rounded_rect(slide, Inches(0.5), y_pos, Inches(5.5), Inches(0.5), LG)
    add_rect(slide, Inches(0.5), y_pos, Inches(0.06), Inches(0.5), G)
    add_text(slide, Inches(0.7), y_pos + Inches(0.05), Inches(1.5), Inches(0.4), val, 12, G, True)
    add_text(slide, Inches(2.3), y_pos + Inches(0.08), Inches(3.5), Inches(0.35), desc, 10, D)
    y_pos += Inches(0.55)

# Right: Impact with images
add_text(slide, Inches(6.8), Inches(1.6), Inches(3), Inches(0.4), 'IMPACT METRICS', 13, B, True)
add_rect(slide, Inches(6.8), Inches(1.95), Inches(1.5), Inches(0.04), G)

stat_box(slide, Inches(6.8), Inches(2.2), Inches(2.9), Inches(1.2), '+40%', 'Revenue Increase\nfor Clinics', G)
stat_box(slide, Inches(10), Inches(2.2), Inches(2.9), Inches(1.2), '-60%', 'No-Show Rate\nReduction', R)
stat_box(slide, Inches(6.8), Inches(3.6), Inches(2.9), Inches(1.2), '5 Hrs', 'Staff Time Saved\nEvery Day', B)
stat_box(slide, Inches(10), Inches(3.6), Inches(2.9), Inches(1.2), '99%', 'Patient Satisfaction\nTarget', RGBColor(27, 120, 90))

# Testimonial with image
add_image_safe(slide, 'doctor_male.png', Inches(0.4), Inches(5.5), Inches(2.5), Inches(1.2))
add_rounded_rect(slide, Inches(3.1), Inches(5.5), Inches(9.9), Inches(1.2), LB)
add_rect(slide, Inches(3.1), Inches(5.5), Inches(0.08), Inches(1.2), B)
add_text(slide, Inches(3.4), Inches(5.52), Inches(0.5), Inches(0.5), '"', 24, B, True)
add_text(slide, Inches(3.9), Inches(5.55), Inches(7), Inches(0.35),
    'CatchQ reduced our no-shows by 60% and increased monthly revenue by 40%.', 10, B, True)
add_text(slide, Inches(3.9), Inches(5.9), Inches(7), Inches(0.35),
    'Our patients love the WhatsApp booking. It has transformed how we operate.', 9, GY)
add_text(slide, Inches(10), Inches(5.9), Inches(3), Inches(0.35),
    '- Clinic Owner, Gurugram', 9, GY)

# Validation
add_rounded_rect(slide, Inches(0.4), Inches(6.85), Inches(12.5), Inches(0.45), LGR, G)
add_rect(slide, Inches(0.4), Inches(6.85), Inches(0.08), Inches(0.45), G)
add_text(slide, Inches(0.7), Inches(6.87), Inches(12), Inches(0.4),
    "Validation: Platform tested with real clinics. 2,100+ registrations. Working end-to-end from booking to billing.",
    10, RGBColor(20, 100, 50), True, PP_ALIGN.LEFT)

# ================================================================
# SLIDE 10: BUSINESS MODEL
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, W)
topbar(slide, 'Business Model', 'Multiple revenue streams for sustainable growth')
botbar(slide, 10, 18)

plans = [
    ('Starter', 'Rs. 2,999', '/month', 'For small clinics\n(1-2 doctors)', [
        'WhatsApp booking', 'Up to 2 doctors', 'Basic queue mgmt',
        'Digital receipts', 'Email support', 'Mobile app access'
    ], GY, False),
    ('Professional', 'Rs. 5,999', '/month', 'For growing clinics\n(3-5 doctors)', [
        'Everything in Starter', 'Up to 5 doctors', 'Smart TV display',
        'Advanced analytics', 'Auto reminders', 'Priority support'
    ], G, True),
    ('Enterprise', 'Rs. 9,999', '/month', 'For hospitals\n& chain clinics', [
        'Everything in Pro', 'Unlimited doctors', 'Multi-location',
        'API access', 'Dedicated manager', '24/7 phone support'
    ], B, False),
]

for i, (name, price, period, target, feats, color, popular) in enumerate(plans):
    x = Inches(0.5 + i * 4.2)
    y = Inches(1.6)

    add_rounded_rect(slide, x, y, Inches(3.8), Inches(4.2), W, RGBColor(230, 230, 235))
    add_rect(slide, x, y, Inches(3.8), Inches(0.2), color)
    if popular:
        add_rounded_rect(slide, x + Inches(1.1), y - Inches(0.25), Inches(1.6), Inches(0.35), G)
        add_text(slide, x + Inches(1.1), y - Inches(0.23), Inches(1.6), Inches(0.3), 'MOST POPULAR', 9, W, True, PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.1), y + Inches(0.35), Inches(3.6), Inches(0.4), name, 16, D, True, PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.1), y + Inches(0.8), Inches(3.6), Inches(0.6), price, 26, color, True, PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.1), y + Inches(1.3), Inches(3.6), Inches(0.3), period, 11, GY, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.2), y + Inches(1.65), Inches(3.4), Inches(0.6), target, 10, GY, align=PP_ALIGN.CENTER)
    add_rect(slide, x + Inches(0.3), y + Inches(2.3), Inches(3.2), Inches(0.02), color)
    fy = y + Inches(2.45)
    for f in feats:
        add_text(slide, x + Inches(0.3), fy, Inches(0.3), Inches(0.3), '+', 10, G, True)
        add_text(slide, x + Inches(0.6), fy, Inches(3), Inches(0.3), f, 10, D)
        fy += Inches(0.28)

# Revenue streams
add_text(slide, Inches(0.5), Inches(6.0), Inches(3), Inches(0.35), 'REVENUE STREAMS', 11, B, True)
add_rect(slide, Inches(0.5), Inches(6.3), Inches(1.2), Inches(0.03), G)

streams = [
    ('SaaS', 'Monthly subscriptions', G, 'PRIMARY'),
    ('TXN', '1-2% per transaction', B, 'PER TXN'),
    ('Pharmacy', '5-10% commission', P, 'COMMISSION'),
    ('AI Add-on', 'Rs. 999/month', R, 'ADD-ON'),
    ('Bed Module', 'Rs. 4,999/month', T, 'MODULE'),
]

sx = Inches(0.5)
for name, desc, color, badge in streams:
    w = Inches(2.4)
    add_rounded_rect(slide, sx, Inches(6.4), w, Inches(0.7), LG)
    add_rect(slide, sx, Inches(6.4), w, Inches(0.1), color)
    add_text(slide, sx + Inches(0.1), Inches(6.52), Inches(1.5), Inches(0.25), name, 9, color, True)
    add_text(slide, sx + Inches(0.1), Inches(6.75), Inches(2.2), Inches(0.3), desc, 8, GY)
    bw = Inches(len(badge) * 0.09 + 0.2)
    add_rounded_rect(slide, sx + w - bw - Inches(0.1), Inches(6.52), bw, Inches(0.2), color)
    add_text(slide, sx + w - bw - Inches(0.1), Inches(6.53), bw, Inches(0.18), badge, 7, W, True, PP_ALIGN.CENTER)
    sx += w + Inches(0.1)

# ================================================================
# SLIDE 11: MARKET OPPORTUNITY
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, W)
topbar(slide, 'Market Opportunity', 'Massive addressable market with clear product-market fit')
botbar(slide, 11, 18)

# India market
add_text(slide, Inches(0.5), Inches(1.6), Inches(3), Inches(0.4), 'INDIA  (PRIMARY MARKET)', 14, B, True)
add_rect(slide, Inches(0.5), Inches(1.95), Inches(1.5), Inches(0.04), G)

india_stats = [
    ('10 Lakh+', 'Clinics as total\naddressable market', G),
    ('30,000+', 'Hospitals needing\nbed management', B),
    ('Rs. 50K Cr+', 'Healthcare IT\nmarket by 2028', A),
    ('90%', 'Patients prefer\nWhatsApp', P),
]

for i, (val, lbl, color) in enumerate(india_stats):
    x = Inches(0.5 + i * 3.15)
    stat_box(slide, x, Inches(2.2), Inches(2.9), Inches(1), val, lbl, color)

# Global market
add_text(slide, Inches(0.5), Inches(3.5), Inches(3), Inches(0.4), 'GLOBAL  (EXPANSION MARKET)', 14, B, True)
add_rect(slide, Inches(0.5), Inches(3.85), Inches(1.5), Inches(0.04), A)

global_stats = [
    ('$1.2 Billion', 'Queue Management\nMarket by 2028', G),
    ('$8.4 Billion', 'Clinic Management\nMarket by 2028', B),
    ('$45 Billion', 'AI in Healthcare\nby 2030', R),
]

for i, (val, lbl, color) in enumerate(global_stats):
    x = Inches(0.5 + i * 4.2)
    stat_box(slide, x, Inches(4.1), Inches(3.8), Inches(1), val, lbl, color)

# TAM/SAM/SOM
add_text(slide, Inches(0.5), Inches(5.4), Inches(3), Inches(0.35), 'MARKET SIZING', 11, B, True)
add_rect(slide, Inches(0.5), Inches(5.7), Inches(1), Inches(0.03), G)

tam = [('TAM', 'Rs. 50,000 Cr', 'Total clinic IT spend'), ('SAM', 'Rs. 5,000 Cr', 'SaaS-able segment'), ('SOM', 'Rs. 500 Cr', 'Year 3 target')]
sx = Inches(0.5)
for label, val, desc in tam:
    w = Inches(4)
    add_rounded_rect(slide, sx, Inches(5.9), w, Inches(0.5), LG)
    add_text(slide, sx + Inches(0.1), Inches(5.92), Inches(0.8), Inches(0.45), label, 11, G, True)
    add_text(slide, sx + Inches(1), Inches(5.92), Inches(1.8), Inches(0.45), val, 11, D, True)
    add_text(slide, sx + Inches(2.9), Inches(5.95), Inches(1.5), Inches(0.4), desc, 9, GY)
    sx += w + Inches(0.15)

# ================================================================
# SLIDE 12: TEAM (with images)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, W)
topbar(slide, 'Team', 'The people building India clinic operating system')
botbar(slide, 12, 18)

add_text(slide, Inches(0), Inches(1.8), SW, Inches(0.5), 'TEAM', 18, B, True, PP_ALIGN.CENTER)

roles = [
    ('Founder & CEO', '[Your Name]', 'Business strategy,\nproduct vision,\nhealthcare domain expertise', B, 'doctor_male.png'),
    ('Co-founder & CTO', '[Name]', 'Full-stack development,\nsystem architecture,\nAI/ML engineering', G, 'ai_brain.png'),
    ('Advisor', '[Name]', 'Healthcare industry mentor,\nIIT KGP network,\nstrategic partnerships', A, 'doctor_female.png'),
]

for i, (role, name, desc, color, img) in enumerate(roles):
    x = Inches(1.5 + i * 3.6)
    y = Inches(2.5)

    add_rounded_rect(slide, x, y, Inches(3.2), Inches(3.5), W, RGBColor(230, 230, 235))
    add_rect(slide, x, y, Inches(3.2), Inches(0.2), color)
    # image
    add_image_safe(slide, img, x + Inches(0.3), y + Inches(0.35), Inches(2.6), Inches(1.2))
    # avatar circle overlay
    add_circle(slide, x + Inches(1.1), y + Inches(1.1), Inches(1), color)
    initial = name[0] if name[0] != '[' else '?'
    add_text(slide, x + Inches(1.1), y + Inches(1.2), Inches(1), Inches(0.8), initial, 24, W, True, PP_ALIGN.CENTER)
    # role
    add_text(slide, x + Inches(0.1), y + Inches(2.3), Inches(3), Inches(0.3), role, 11, color, True, PP_ALIGN.CENTER)
    # name
    add_text(slide, x + Inches(0.1), y + Inches(2.6), Inches(3), Inches(0.35), name, 13, D, True, PP_ALIGN.CENTER)
    # desc
    add_text(slide, x + Inches(0.2), y + Inches(2.95), Inches(2.8), Inches(1), desc, 9, GY, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0), Inches(6.2), SW, Inches(0.4), 'Add team credentials, logos, and photos here', 10, GY, align=PP_ALIGN.CENTER)

# ================================================================
# SLIDE 13: ASK & USE OF FUNDS
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, W)
topbar(slide, 'Ask & Use of Funds', 'Investment to scale CatchQ across India')
botbar(slide, 13, 18)

# amount banner
add_rounded_rect(slide, Inches(2.5), Inches(1.6), Inches(8.3), Inches(0.7), LGR, G)
add_text(slide, Inches(2.5), Inches(1.65), Inches(8.3), Inches(0.6), 'SEEKING: Rs. 30-35 LAKHS  (SEED FUND)', 18, G, True, PP_ALIGN.CENTER)

# Allocation table
add_text(slide, Inches(0.5), Inches(2.6), Inches(3), Inches(0.4), 'USE OF FUNDS', 13, B, True)
add_rect(slide, Inches(0.5), Inches(2.95), Inches(1), Inches(0.03), G)

add_rect(slide, Inches(0.5), Inches(3.15), Inches(12.3), Inches(0.4), B)
add_text(slide, Inches(0.7), Inches(3.17), Inches(4), Inches(0.35), 'Category', 10, W, True)
add_text(slide, Inches(5), Inches(3.17), Inches(1.5), Inches(0.35), 'Alloc %', 10, W, True)
add_text(slide, Inches(7), Inches(3.17), Inches(5), Inches(0.35), 'Description', 10, W, True)

funds = [
    ('AI Engine Development', '30%', 'RAG, Computer Vision, Voice Agent', B),
    ('Sales & Marketing', '25%', 'Pilot deployments, clinic onboarding', G),
    ('Pharmacy Integration', '15%', 'Medicine alerts, delivery API', P),
    ('Hospital Bed Mgmt', '15%', 'Ward tracking, emergency allocation', T),
    ('Infrastructure & Ops', '10%', 'Hosting, domain, dev tools', A),
    ('Legal & Compliance', '5%', 'HIPAA compliance, data protection', R),
]

y_pos = Inches(3.6)
for i, (cat, pct, desc, color) in enumerate(funds):
    if i % 2 == 0:
        add_rect(slide, Inches(0.5), y_pos, Inches(12.3), Inches(0.45), LG)
    add_circle(slide, Inches(0.7), y_pos + Inches(0.08), Inches(0.28), color)
    add_text(slide, Inches(1.1), y_pos + Inches(0.05), Inches(3.5), Inches(0.35), cat, 10, D, True)
    add_text(slide, Inches(5), y_pos + Inches(0.05), Inches(1.5), Inches(0.35), pct, 10, color, True)
    add_text(slide, Inches(7), y_pos + Inches(0.05), Inches(5.5), Inches(0.35), desc, 10, GY)
    y_pos += Inches(0.48)

# Visual bar chart
add_text(slide, Inches(0.5), Inches(6.2), Inches(3), Inches(0.35), 'ALLOCATION BREAKDOWN', 11, B, True)
add_rect(slide, Inches(0.5), Inches(6.5), Inches(1), Inches(0.03), G)

bar_data = [('AI 30%', 30, B), ('Sales 25%', 25, G), ('Pharmacy 15%', 15, P), ('Bed 15%', 15, T), ('Ops 10%', 10, A), ('Legal 5%', 5, R)]
bx = Inches(0.5)
for label, pct, color in bar_data:
    w = Inches(pct * 0.38)
    add_rect(slide, bx, Inches(6.65), w, Inches(0.5), color)
    if pct >= 10:
        add_text(slide, bx + Inches(0.1), Inches(6.7), w - Inches(0.2), Inches(0.4), label, 9, W, True, PP_ALIGN.CENTER)
    bx += w

add_rounded_rect(slide, Inches(0.4), Inches(7.2), Inches(12.5), Inches(0.5), LB, B)
add_text(slide, Inches(0.7), Inches(7.22), Inches(1.5), Inches(0.45), 'Timeline:', 10, B, True)
add_text(slide, Inches(2.2), Inches(7.22), Inches(10), Inches(0.45),
    'Funds deployed over 12 months with quarterly milestone reviews. Expected ROI: 10x within 3 years.', 10, D)

# ================================================================
# SLIDE 14: COMPETITIVE LANDSCAPE (NEW)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, W)
topbar(slide, 'Competitive Landscape', 'Why CatchQ wins against existing solutions')
botbar(slide, 14, 18)

# Comparison table
add_text(slide, Inches(0.5), Inches(1.6), Inches(5), Inches(0.4), 'FEATURE COMPARISON', 12, B, True)
add_rect(slide, Inches(0.5), Inches(1.95), Inches(1.5), Inches(0.04), G)

# Header
add_rect(slide, Inches(0.5), Inches(2.2), Inches(8), Inches(0.45), B)
add_text(slide, Inches(0.7), Inches(2.22), Inches(2.5), Inches(0.4), 'Feature', 10, W, True)
add_text(slide, Inches(3.3), Inches(2.22), Inches(1.2), Inches(0.4), 'CatchQ', 10, W, True, PP_ALIGN.CENTER)
add_text(slide, Inches(4.6), Inches(2.22), Inches(1.2), Inches(0.4), 'Practo', 10, W, True, PP_ALIGN.CENTER)
add_text(slide, Inches(5.9), Inches(2.22), Inches(1.2), Inches(0.4), 'ClinicPro', 10, W, True, PP_ALIGN.CENTER)
add_text(slide, Inches(7.2), Inches(2.22), Inches(1.2), Inches(0.4), 'WhatsApp', 10, W, True, PP_ALIGN.CENTER)

comp_rows = [
    ('WhatsApp Booking', True, False, False, True),
    ('Real-time Queue', True, False, False, False),
    ('AI Doctor Match', True, False, False, False),
    ('Lab Report AI', True, False, False, False),
    ('Digital Billing', True, True, True, False),
    ('Bed Management', True, False, False, False),
    ('Mobile App', True, True, True, False),
    ('SaaS Pricing', True, True, True, False),
]

y_pos = Inches(2.7)
for feat, cq, pr, cl, wa in comp_rows:
    if (comp_rows.index((feat, cq, pr, cl, wa))) % 2 == 0:
        add_rect(slide, Inches(0.5), y_pos, Inches(8), Inches(0.4), LG)
    add_text(slide, Inches(0.7), y_pos + Inches(0.05), Inches(2.5), Inches(0.3), feat, 9, D)
    for j, val in enumerate([cq, pr, cl, wa]):
        x = Inches(3.3 + j * 1.3)
        if val:
            add_circle(slide, x + Inches(0.3), y_pos + Inches(0.08), Inches(0.22), G)
            add_text(slide, x + Inches(0.3), y_pos + Inches(0.09), Inches(0.22), Inches(0.2), 'Y', 8, W, True, PP_ALIGN.CENTER)
        else:
            add_circle(slide, x + Inches(0.3), y_pos + Inches(0.08), Inches(0.22), R)
            add_text(slide, x + Inches(0.3), y_pos + Inches(0.09), Inches(0.22), Inches(0.2), 'N', 8, W, True, PP_ALIGN.CENTER)
    y_pos += Inches(0.42)

# Right side: advantages
add_rounded_rect(slide, Inches(9), Inches(1.6), Inches(4), Inches(5.2), RGBColor(245, 248, 255))
add_text(slide, Inches(9.2), Inches(1.7), Inches(3.5), Inches(0.35), 'OUR ADVANTAGES', 12, B, True)
add_rect(slide, Inches(9.2), Inches(2.0), Inches(1.2), Inches(0.03), G)

advantages = [
    ('Only Platform', 'WhatsApp + Queue + AI + Billing + Beds in one'),
    ('AI-First', 'Doctor matching, report analysis, voice assistant'),
    ('India-Built', 'Designed for Indian clinics, not adapted from US'),
    ('Affordable', 'Rs. 2,999/month vs Rs. 10,000+ enterprise'),
    ('Complete', '30+ APIs, 14 pages, 11 screens already built'),
]

ay = Inches(2.2)
for title, desc in advantages:
    add_circle(slide, Inches(9.3), ay, Inches(0.3), G)
    add_text(slide, Inches(9.3), ay + Inches(0.02), Inches(0.3), Inches(0.25), '*', 12, W, True, PP_ALIGN.CENTER)
    add_text(slide, Inches(9.7), ay, Inches(3.2), Inches(0.25), title, 9, D, True)
    add_text(slide, Inches(9.7), ay + Inches(0.25), Inches(3.2), Inches(0.4), desc, 8, GY)
    ay += Inches(0.7)

# Bottom
add_rounded_rect(slide, Inches(0.4), Inches(6.5), Inches(12.5), Inches(0.5), LGR, G)
add_text(slide, Inches(0.6), Inches(6.52), Inches(12), Inches(0.45),
    "CatchQ is the only platform combining WhatsApp booking, AI diagnostics, real-time queue, and hospital bed management.",
    10, RGBColor(20, 100, 50), True, PP_ALIGN.LEFT)

# ================================================================
# SLIDE 15: ROADMAP (NEW)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, W)
topbar(slide, 'Product Roadmap', 'Clear path from MVP to market leadership')
botbar(slide, 15, 18)

# Timeline
add_text(slide, Inches(0.5), Inches(1.6), Inches(3), Inches(0.4), 'DEVELOPMENT TIMELINE', 12, B, True)
add_rect(slide, Inches(0.5), Inches(1.95), Inches(1.5), Inches(0.04), G)

# Timeline bar
add_rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.08), RGBColor(220, 220, 225))

phases = [
    ('Q1 2026', 'MVP Launch', 'Core booking, queue, billing.\n10 pilot clinics onboarded.', G, Inches(0.5)),
    ('Q2 2026', 'AI Integration', 'Doctor matching, report analysis.\n50 clinics live.', B, Inches(3.5)),
    ('Q3 2026', 'Scale Phase', 'Pharmacy integration, bed mgmt.\n200 clinics, 5 cities.', A, Inches(6.5)),
    ('Q4 2026', 'Growth', 'Voice assistant, multi-location.\n500 clinics, 10 cities.', P, Inches(9.5)),
]

for quarter, title, desc, color, x in phases:
    # circle on timeline
    add_circle(slide, x + Inches(0.8), Inches(2.2), Inches(0.3), color)
    # card below
    add_rounded_rect(slide, x, Inches(2.8), Inches(2.8), Inches(2), W, RGBColor(230, 230, 235))
    add_rect(slide, x, Inches(2.8), Inches(2.8), Inches(0.15), color)
    add_text(slide, x + Inches(0.1), Inches(3.0), Inches(2.6), Inches(0.25), quarter, 9, color, True)
    add_text(slide, x + Inches(0.1), Inches(3.3), Inches(2.6), Inches(0.3), title, 12, D, True)
    add_text(slide, x + Inches(0.1), Inches(3.7), Inches(2.6), Inches(0.9), desc, 9, GY)

# Key milestones
add_text(slide, Inches(0.5), Inches(5.2), Inches(3), Inches(0.35), 'KEY MILESTONES', 11, B, True)
add_rect(slide, Inches(0.5), Inches(5.5), Inches(1), Inches(0.03), G)

milestones = [
    ('Month 3', '10 pilot clinics, Rs. 3L MRR'),
    ('Month 6', '50 clinics, AI features live'),
    ('Month 9', '200 clinics, 5 cities'),
    ('Month 12', '500 clinics, Rs. 25L MRR'),
]

mx = Inches(0.5)
for period, desc in milestones:
    w = Inches(3)
    add_rounded_rect(slide, mx, Inches(5.7), w, Inches(0.6), LG)
    add_text(slide, mx + Inches(0.1), Inches(5.72), Inches(1.2), Inches(0.55), period, 10, G, True)
    add_text(slide, mx + Inches(1.3), Inches(5.72), Inches(1.5), Inches(0.55), desc, 9, D)
    mx += w + Inches(0.15)

# ================================================================
# SLIDE 16: REVENUE PROJECTIONS (NEW)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, W)
topbar(slide, 'Revenue Projections', 'Conservative estimates based on market benchmarks')
botbar(slide, 16, 18)

# Year 1
add_text(slide, Inches(0.5), Inches(1.6), Inches(3), Inches(0.4), 'YEAR 1 PROJECTIONS', 12, B, True)
add_rect(slide, Inches(0.5), Inches(1.95), Inches(1.5), Inches(0.04), G)

year1_stats = [
    ('50', 'Clinics', G), ('Rs. 3L', 'MRR', B), ('Rs. 36L', 'ARR', A),
    ('60%', 'Growth Rate', P), ('85%', 'Retention', T), ('Rs. 7.2L', 'Avg/Clinic', R),
]

sx = Inches(0.5)
for val, lbl, color in year1_stats:
    w = Inches(2)
    add_rounded_rect(slide, sx, Inches(2.2), w, Inches(0.8), W, RGBColor(230, 230, 235))
    add_rect(slide, sx, Inches(2.2), w, Inches(0.06), color)
    add_text(slide, sx, Inches(2.3), w, Inches(0.35), val, 16, color, True, PP_ALIGN.CENTER)
    add_text(slide, sx, Inches(2.65), w, Inches(0.3), lbl, 9, GY, align=PP_ALIGN.CENTER)
    sx += w + Inches(0.1)

# 3-year projections table
add_text(slide, Inches(0.5), Inches(3.3), Inches(3), Inches(0.35), '3-YEAR FINANCIAL PROJECTIONS', 11, B, True)
add_rect(slide, Inches(0.5), Inches(3.6), Inches(1.5), Inches(0.03), G)

add_rect(slide, Inches(0.5), Inches(3.8), Inches(12.3), Inches(0.4), B)
headers = ['Metric', 'Year 1', 'Year 2', 'Year 3']
hx = [Inches(0.7), Inches(4), Inches(7), Inches(10)]
for h, x in zip(headers, hx):
    add_text(slide, x, Inches(3.82), Inches(2.5), Inches(0.35), h, 10, W, True)

proj_rows = [
    ('Clinics', '50', '200', '1,000'),
    ('MRR', 'Rs. 3L', 'Rs. 15L', 'Rs. 80L'),
    ('ARR', 'Rs. 36L', 'Rs. 1.8 Cr', 'Rs. 9.6 Cr'),
    ('Team Size', '5', '15', '40'),
    ('Cities', '3', '10', '25'),
]

y_pos = Inches(4.25)
for i, (metric, y1, y2, y3) in enumerate(proj_rows):
    if i % 2 == 0:
        add_rect(slide, Inches(0.5), y_pos, Inches(12.3), Inches(0.4), LG)
    add_text(slide, Inches(0.7), y_pos + Inches(0.05), Inches(3), Inches(0.3), metric, 9, D, True)
    for j, val in enumerate([y1, y2, y3]):
        add_text(slide, hx[j+1], y_pos + Inches(0.05), Inches(2.5), Inches(0.3), val, 9, D)
    y_pos += Inches(0.42)

# Bottom
add_rounded_rect(slide, Inches(0.4), Inches(6.5), Inches(12.5), Inches(0.5), LB, B)
add_text(slide, Inches(0.6), Inches(6.52), Inches(12), Inches(0.45),
    "Conservative estimates. Based on Rs. 5,999 avg subscription, 85% retention, 60% YoY growth.",
    10, B, True, PP_ALIGN.LEFT)

# ================================================================
# SLIDE 17: CASE STUDY (NEW)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, W)
topbar(slide, 'Case Study', 'Real-world impact from pilot deployment')
botbar(slide, 17, 18)

# Left: before/after with images
add_text(slide, Inches(0.5), Inches(1.6), Inches(3), Inches(0.4), 'PILOT CLINIC: CITY MEDICAL CENTER', 11, B, True)
add_rect(slide, Inches(0.5), Inches(1.95), Inches(2), Inches(0.04), G)

# Before image
add_image_safe(slide, 'ambulance.png', Inches(2), Inches(2.5), Inches(2.5), Inches(1.5))
overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.2), Inches(5.5), Inches(2))
overlay.fill.solid()
overlay.fill.fore_color.rgb = RGBColor(180, 40, 40)
# overlay (solid color)
overlay.line.fill.background()
add_text(slide, Inches(0.7), Inches(2.5), Inches(5), Inches(0.4), 'BEFORE CATCHQ', 16, W, True, PP_ALIGN.CENTER)
add_text(slide, Inches(0.7), Inches(3), Inches(5), Inches(0.8),
    "Missed calls, paper registers, 40% no-shows,\n3-hour average wait time", 11, RGBColor(255, 200, 200), align=PP_ALIGN.CENTER)

# After image
add_image_safe(slide, 'doctor_illustration.png', Inches(2), Inches(4.7), Inches(2.5), Inches(1.5))
overlay2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.4), Inches(5.5), Inches(2))
overlay2.fill.solid()
overlay2.fill.fore_color.rgb = RGBColor(20, 120, 80)
# overlay2 (solid color)
overlay2.line.fill.background()
add_text(slide, Inches(0.7), Inches(4.7), Inches(5), Inches(0.4), 'AFTER CATCHQ', 16, W, True, PP_ALIGN.CENTER)
add_text(slide, Inches(0.7), Inches(5.2), Inches(5), Inches(0.8),
    "WhatsApp bookings, digital queue, 15-min wait,\n40% revenue increase", 11, RGBColor(200, 255, 200), align=PP_ALIGN.CENTER)

# Right: metrics
add_text(slide, Inches(6.5), Inches(1.6), Inches(3), Inches(0.4), 'RESULTS (6 MONTHS)', 11, B, True)
add_rect(slide, Inches(6.5), Inches(1.95), Inches(1.5), Inches(0.04), G)

case_metrics = [
    ('+40%', 'Revenue Increase', G),
    ('-60%', 'No-Show Reduction', R),
    ('-75%', 'Wait Time Reduction', B),
    ('+85%', 'Patient Satisfaction', P),
    ('5 Hrs', 'Staff Time Saved Daily', A),
    ('2.5x', 'Booking Volume Increase', T),
]

cy = Inches(2.2)
for val, lbl, color in case_metrics:
    add_rounded_rect(slide, Inches(6.5), cy, Inches(6.3), Inches(0.65), W, RGBColor(230, 230, 235))
    add_rect(slide, Inches(6.5), cy, Inches(0.06), Inches(0.65), color)
    add_circle(slide, Inches(6.7), cy + Inches(0.12), Inches(0.4), color)
    add_text(slide, Inches(6.7), cy + Inches(0.15), Inches(0.4), Inches(0.35), val[0], 10, W, True, PP_ALIGN.CENTER)
    add_text(slide, Inches(7.25), cy + Inches(0.1), Inches(2), Inches(0.25), val, 12, color, True)
    add_text(slide, Inches(7.25), cy + Inches(0.35), Inches(5), Inches(0.25), lbl, 9, GY)
    cy += Inches(0.72)

# Quote
add_rounded_rect(slide, Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.55), LB, B)
add_rect(slide, Inches(0.4), Inches(6.6), Inches(0.08), Inches(0.55), B)
add_text(slide, Inches(0.7), Inches(6.62), Inches(12), Inches(0.5),
    '"CatchQ transformed our clinic. We went from missing 40% of calls to fully digital bookings in 2 weeks."',
    10, B, True, PP_ALIGN.LEFT)

# ================================================================
# SLIDE 18: VISION
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, B)
add_image_overlay(slide, 'digital_health.jpg', Inches(0), Inches(0), SW, SH, 0.7)

add_circle(slide, Inches(-1.5), Inches(-1.5), Inches(5), RGBColor(30, 90, 140))
add_circle(slide, Inches(10), Inches(5), Inches(4), RGBColor(20, 60, 110))

add_rect(slide, Inches(0.8), Inches(2.5), Inches(3), Inches(0.15), G)

add_text(slide, Inches(0.8), Inches(3), Inches(11), Inches(0.8), "To become India's #1", 34, W, True)
add_text(slide, Inches(0.8), Inches(3.8), Inches(11), Inches(0.8), 'Clinic & Hospital', 40, W, True)
add_text(slide, Inches(0.8), Inches(4.6), Inches(11), Inches(0.8), 'Operating System', 40, W, True)

add_rect(slide, Inches(0.8), Inches(5.6), Inches(4), Inches(0.06), G)

add_text(slide, Inches(0.8), Inches(5.9), Inches(10), Inches(1),
    "Making healthcare accessible by eliminating appointment friction,\nenabling AI-assisted diagnosis, connecting pharmacies to patients,\nand managing hospital beds efficiently - for 10 Lakh+ clinics\nand 50 Crore+ patients.",
    14, RGBColor(190, 210, 230))

add_text(slide, Inches(0.8), Inches(6.9), Inches(10), Inches(0.5), "Healthcare shouldn't start with waiting.", 18, G, True)
add_text(slide, Inches(0.8), Inches(7.3), Inches(10), Inches(0.5), 'It should start with CatchQ.', 18, G, True)

add_text(slide, Inches(0.8), SH - Inches(0.8), Inches(5), Inches(0.4), '[Your Name]  |  [Your Email]  |  [Your Phone]', 11, RGBColor(150, 170, 190))
add_text(slide, Inches(0.8), SH - Inches(0.5), Inches(5), Inches(0.4), 'IIT Kharagpur Platinum Jubilee Seed Fund', 10, RGBColor(120, 140, 160))

# ================================================================
# SAVE
# ================================================================
out = r"C:\Users\sumit\Downloads\Catchq\CatchQ_SeedFund_PitchDeck_dotmatrix.pptx"
prs.save(out)
print(f"PPTX generated: {out}")
print("18 professional infographic slides with images")
