from pptx import Presentation
import copy
from lxml import etree

MAIN = r"C:\Users\sumit\Downloads\Catchq\CatchQ_SeedFund_PitchDeck.pptx"
NEW_SLIDES = r"C:\Users\sumit\Downloads\Catchq\CatchQ_Queue_TAM_Slides.pptx"
OUT = r"C:\Users\sumit\Downloads\Catchq\CatchQ_SeedFund_PitchDeck_final.pptx"

prs = Presentation(MAIN)
new_prs = Presentation(NEW_SLIDES)

print(f"Main deck: {len(prs.slides)} slides")
print(f"New slides: {len(new_prs.slides)} slides")

# Add all 5 new slides to the main deck
for i, slide_layout in enumerate(new_prs.slide_layouts):
    pass  # just counting

for slide in new_prs.slides:
    # Copy the slide by appending its XML
    slide_layout = prs.slide_layouts[6]  # blank layout
    new_slide = prs.slides.add_slide(slide_layout)
    
    # Remove default placeholders
    for shape in list(new_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)
    
    # Copy all shapes from source slide
    for shape in slide.shapes:
        el = copy.deepcopy(shape._element)
        new_slide.shapes._spTree.append(el)

prs.save(OUT)
print(f"Final deck: {len(prs.slides)} slides")
print(f"Saved: {OUT}")
